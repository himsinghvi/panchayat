"""Generate Panchaayat investor / VC / B2B pitch deck."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
PRIMARY = RGBColor(0x25, 0x63, 0xEB)      # #2563eb
PRIMARY_DARK = RGBColor(0x1D, 0x4E, 0xD8)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)       # amber
DARK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)

OUTPUT = Path(__file__).resolve().parent.parent / "Panchaayat-Pitch-Deck.pptx"


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, text="Panchaayat — Confidential"):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = SLATE


def add_title_bar(slide, title: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xDB, 0xEA, 0xFE)


def add_bullets(slide, items, left=0.6, top=1.45, width=8.8, height=5.5, font_size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
        p.level = 0


def add_two_column(slide, left_title, left_items, right_title, right_items):
    # Left column
    lt = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(4.2), Inches(0.4))
    lt.text_frame.paragraphs[0].text = left_title
    lt.text_frame.paragraphs[0].font.bold = True
    lt.text_frame.paragraphs[0].font.size = Pt(18)
    lt.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    add_bullets(slide, left_items, left=0.6, top=1.85, width=4.2, height=4.8, font_size=14)

    # Right column
    rt = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.2), Inches(0.4))
    rt.text_frame.paragraphs[0].text = right_title
    rt.text_frame.paragraphs[0].font.bold = True
    rt.text_frame.paragraphs[0].font.size = Pt(18)
    rt.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    add_bullets(slide, right_items, left=5.2, top=1.85, width=4.2, height=4.8, font_size=14)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY_DARK)
    # Accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.8), Inches(10), Inches(0.08))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.5), Inches(1.2))
    p = t.text_frame.paragraphs[0]
    p.text = "Panchaayat"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE

    s = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(8.5), Inches(1.0))
    sp = s.text_frame.paragraphs[0]
    sp.text = "India's public record for consumer experiences — and a case-management system SMEs can trust"
    sp.font.size = Pt(22)
    sp.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

    tag = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(8), Inches(0.5))
    tp = tag.text_frame.paragraphs[0]
    tp.text = "Investor & Partnership Pitch  |  2026"
    tp.font.size = Pt(14)
    tp.font.color.rgb = ACCENT

    add_footer(slide, "github.com/himsinghvi/panchayat")


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "The Problem", "Consumers are heard. Brands are blind. Nobody closes the loop.")
    add_bullets(slide, [
        "Reviews live on Google, complaints on NCH, grievances on social — all fragmented, none verified",
        "Brands can reply anywhere, but no platform tracks whether the actual customer got resolution",
        "SMEs run support through shared inboxes — no case IDs, no SLAs, no visibility into open issues",
        "Consumers have no single place to see a brand's real resolution track record (not just star ratings)",
        "Existing platforms let businesses mark issues 'resolved' without consumer confirmation",
    ])
    add_footer(slide)


def slide_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Market Opportunity", "India — massive consumer economy, rising digital grievance culture")
    add_two_column(slide,
        "Consumer side",
        [
            "63M+ MSMEs in India; most lack formal helpdesk software",
            "Rising consumer awareness (NCH, E-Jagriti, social escalation)",
            "Google reviews ≠ accountability; MouthShut/complaint sites lack resolution loops",
            "Hyper-local commerce (D2C, furniture, appliances) = high post-purchase friction",
        ],
        "Business side",
        [
            "Global customer service software market: $50B+ and growing",
            "India SaaS for SMB support is underserved vs. Zendesk/Freshdesk price points",
            "Brands need reputation management + structured case management in one place",
            "SEO on '[Brand] complaints' creates organic demand flywheel",
        ])
    add_footer(slide)


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "The Solution — Panchaayat", "One platform. Two engines. One trust layer.")
    add_bullets(slide, [
        "PUBLIC LAYER — Consumers share experiences, discuss publicly, escalate when brands ignore",
        "PRIVATE LAYER — SMEs get branded case-management pages (email-to-case, SLA dashboards)",
        "TRUST LAYER — Only the original consumer confirms resolution; full history is immutable",
        "AI LAYER — Complaint drafting, quality checks, smart search, contextual ad targeting",
        "Tagline: Make it visible. Get it resolved. Verified by you.",
    ], font_size=17)
    add_footer(slide)


def slide_how_it_works(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "How It Works", "Consumer journey — transparent from post to resolution")
    steps = [
        "1. SHARE — Post review/complaint with brand, location, evidence (guest or registered)",
        "2. DISCUSS — Community comments, 'Me Too', AI summary; thread builds public pressure",
        "3. BRAND RESPONDS — Verified rep replies officially on the public record",
        "4. RESOLVE — Brand proposes refund/replacement/repair with proof",
        "5. YOU CONFIRM — Only the original consumer closes the case (reject & reopen if insufficient)",
        "6. TRACK RECORD — Brand profile shows resolution rate, response time, category trends",
    ]
    add_bullets(slide, steps, font_size=15)
    add_footer(slide)


def slide_differentiator(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Key Differentiator", "The resolution loop no one else enforces")
    add_two_column(slide,
        "What others do",
        [
            "Star ratings with no accountability",
            "Brands self-mark as 'resolved'",
            "Complaints disappear into forums",
            "No link between public reputation & private support",
        ],
        "What Panchaayat does",
        [
            "Consumer-confirmed resolution only",
            "Immutable timeline — every status change visible",
            "Public + private cases on same platform",
            "Brands cannot pay to remove legitimate complaints",
            "AI assists without judging guilt",
        ])
    add_footer(slide)


def slide_product(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Product — What's Built (MVP Live)", "Full-stack platform deployed and demo-ready")
    add_bullets(slide, [
        "Consumer: complaint wizard, AI draft & quality check, feed, smart search, dashboards",
        "Brand: verified profiles, inbox, KPI dashboard, official replies, resolution proposals",
        "Admin: contextual ads, AI targeting, moderation-ready architecture",
        "B2B pages: pricing tiers, API docs, hot-query integrations (DND, refund status, etc.)",
        "Tech: FastAPI + React, JWT auth, Azure OpenAI with fallbacks, dark/light themes",
        "Demo: 8 personas, 6 brands, seeded complaints — login at /login (password: demo123)",
    ], font_size=15)
    add_footer(slide)


def slide_personas(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Target Users", "Multi-sided platform with clear personas")
    add_two_column(slide,
        "Consumers & Community",
        [
            "Complainants — post, track, confirm resolution",
            "Community members — discuss, Me Too, support",
            "Guests — post immediately (lower weight score)",
            "Verified buyers — higher credibility & brand priority",
        ],
        "Business & Platform",
        [
            "Brand reps — respond, propose resolutions, analytics",
            "SME owners — branded support page, email-to-case",
            "Moderators — abuse reports, verification",
            "Enterprise — API, webhooks, white-label, call centre",
        ])
    add_footer(slide)


def slide_b2b(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "B2B — For Business", "Turn your inbox into a real case management system")
    add_bullets(slide, [
        "Problem: SMEs run support via shared email — queries lost, no SLA, no ownership",
        "Solution: Branded page at panchaayat.in/yourbrand + 2-way email sync + live dashboard",
        "Example: Chaitanya Furnishings (40-person D2C) — 300 emails/month → tracked cases in 30 days",
        "Private cases by default; public escalation only if unresolved",
        "Features: CSAT/NPS, social escalation scraping, SLA rules, secure attachments",
        "30-day money-back guarantee — no lock-in contracts",
    ], font_size=15)
    add_footer(slide)


def slide_pricing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Business Model — SaaS Pricing", "Recurring revenue from SME & enterprise tiers")
    # Three pricing cards as text boxes
    plans = [
        ("Startup", "₹1,999/mo", "150 cases, 1 seat, branded page, email+SMS"),
        ("Scale-up", "₹3,999/mo", "Unlimited cases, 10 seats, email sync, CSAT, social tagging"),
        ("Enterprise", "₹5,999/mo", "API, AI triage, white-label, webhooks, unlimited seats"),
    ]
    for i, (name, price, feats) in enumerate(plans):
        left = 0.5 + i * 3.15
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(2.9), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE if i != 1 else RGBColor(0xEF, 0xF6, 0xFF)
        card.line.color.rgb = PRIMARY if i == 1 else RGBColor(0xE2, 0xE8, 0xF0)

        tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(1.7), Inches(2.5), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, (text, bold, size, color) in enumerate([
            (name, True, 20, PRIMARY),
            (price, True, 24, DARK),
            (feats, False, 12, SLATE),
        ]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = text
            p.font.bold = bold
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(8)

    note = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(8.8), Inches(0.5))
    note.text_frame.paragraphs[0].text = "+ Custom / Call Centre plan for large corporates (outsourced voice, dedicated CSM, custom SLAs)"
    note.text_frame.paragraphs[0].font.size = Pt(12)
    note.text_frame.paragraphs[0].font.color.rgb = SLATE
    add_footer(slide)


def slide_revenue(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Revenue Streams", "Trust-first monetization — ads never bias complaint visibility")
    add_bullets(slide, [
        "SaaS subscriptions — Startup / Scale-up / Enterprise (primary near-term revenue)",
        "Enterprise API & white-label — per-seat or usage-based for large brands",
        "Contextual ads — useful services (legal aid, warranty, refund trackers) matched to complaint context",
        "Future: verified brand analytics, priority response tools, researcher API (rate-limited)",
        "Explicitly avoided: pay-to-remove complaints, biased ranking — kills core trust proposition",
    ])
    add_footer(slide)


def slide_competition(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Competitive Landscape", "Fragmented incumbents — no one owns the full loop")
    add_bullets(slide, [
        "Google Reviews / Amazon — ratings only, no resolution workflow, no consumer confirmation",
        "MouthShut / ConsumerComplaints.in — complaint listing, weak brand engagement & resolution tracking",
        "NCH / E-Jagriti — government escalation, not public discussion or brand reputation layer",
        "Zendesk / Freshdesk — enterprise helpdesk, too expensive & complex for Indian SMEs",
        "Panchaayat wedge: public reputation + private case management + consumer-verified resolution",
    ], font_size=15)
    add_footer(slide)


def slide_gtm(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Go-to-Market Strategy", "Hyper-local density → national scale")
    add_bullets(slide, [
        "Phase 1 — Seed hyper-local (e.g., Pune/Hyderabad): complaints + brand profiles with demo data",
        "Phase 2 — SEO flywheel: rank for '[Brand] complaints [City]' — organic consumer & brand traffic",
        "Phase 3 — SME outbound: D2C brands, furniture, appliances — 'replace your support@ email link'",
        "Phase 4 — Enterprise API partnerships: banks, insurers, telcos (hot queries: DND, claim status)",
        "Community growth: Me Too, trending feed, recently resolved — social proof drives trust",
        "14-day free trial for business accounts — low friction conversion",
    ], font_size=15)
    add_footer(slide)


def slide_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Technology & Moat", "Built for India, extensible globally")
    add_two_column(slide,
        "Stack (live MVP)",
        [
            "FastAPI + React + SQLite/Postgres",
            "Azure OpenAI + rule-based fallbacks",
            "JWT auth, REST API, webhooks (Enterprise)",
            "Smart search: synonyms, entity extraction",
            "Deployed: Vercel + single-app architecture",
        ],
        "Defensibility",
        [
            "Resolution history data — hard to replicate",
            "Brand resolution scores — reputation graph",
            "Consumer-confirmed outcomes — trust signal",
            "SME email-to-case integration — workflow lock-in",
            "India-specific: NCH links, local language roadmap",
        ])
    add_footer(slide)


def slide_metrics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Success Metrics & Traction", "What we measure and where we are today")
    add_two_column(slide,
        "North-star KPIs",
        [
            "% complaints with brand response < 7 days",
            "% consumer-confirmed resolutions (not brand-claimed)",
            "Median time-to-resolution",
            "Monthly active complainants & brand responders",
            "Organic search traffic to brand/complaint pages",
            "SME trial → paid conversion rate",
        ],
        "Current status",
        [
            "MVP live with full consumer + brand flows",
            "8 demo personas, 6 seeded brands",
            "AI features operational (draft, search, ads)",
            "B2B pricing & API docs published",
            "Open source on GitHub (himsinghvi/panchayat)",
            "Seeking: pilot SMEs + seed funding",
        ])
    add_footer(slide)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Roadmap", "12–18 month product & GTM plan")
    add_bullets(slide, [
        "Q2 2026 — Pilot 10 SME brands (Pune/Bengaluru); Postgres migration; mobile-responsive PWA",
        "Q3 2026 — Email-to-case sync; Hindi/Marathi UI; NCH/E-Jagriti escalation integration",
        "Q4 2026 — Social escalation scraping; CSAT/NPS; 100 paying SME accounts",
        "Q1 2027 — Enterprise API GA; white-label; AI triage chatbot; 5 enterprise logos",
        "Q2 2027 — Mobile apps; verified-purchase flow; regional expansion (Tier-2 cities)",
        "Ongoing — Legal/ToS review; moderation scale; credibility scoring algorithm",
    ], font_size=14)
    add_footer(slide)


def slide_risks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Risks & Mitigations", "Honest assessment for investors")
    add_bullets(slide, [
        "Defamation liability → Clear ToS, evidence requirements, takedown/appeal process, legal review",
        "Fake reviews / astroturfing → Verified purchase signals, moderation queue, consumer-only resolution confirm",
        "Cold start → Hyper-local seeding, SEO pages, SME trials, public data aggregation (unclaimed brands)",
        "Brand incentive → Public pages rank on Google; ignoring complaints costs reputation",
        "SME churn → Email sync lock-in, resolution analytics ROI, 30-day money-back reduces risk",
    ], font_size=15)
    add_footer(slide)


def slide_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)
    add_title_bar(slide, "Team", "[Founder name & background — customize before presenting]")
    add_bullets(slide, [
        "Founder — Product & engineering (built full MVP solo: backend, frontend, AI, deployment)",
        "Advisors needed — Consumer law (India), SME GTM, enterprise sales",
        "Hiring plan (post-seed) — Full-stack engineer, community ops, SME sales lead",
        "Location — India (Pune / remote-first)",
        "",
        "→ Replace this slide with actual team bios, photos, and relevant experience before investor meetings.",
    ], font_size=16)
    add_footer(slide)


def slide_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY_DARK)
    add_title_bar(slide, "The Ask", "Seed round to scale GTM and product")
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    items = [
        ("Raising", "₹[X] Cr seed round  (customize amount)"),
        ("Use of funds", "40% GTM & SME sales  |  30% engineering  |  20% ops & legal  |  10% buffer"),
        ("Milestones (12 mo)", "100 paying SMEs  |  10K monthly active users  |  Enterprise pilot"),
        ("For VCs", "Category-defining trust infrastructure for India's consumer economy"),
        ("For business partners", "14-day free trial — branded support page live in < 1 day"),
        ("For strategic investors", "Data on resolution patterns, brand responsiveness, category trends"),
    ]
    for i, (label, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{label}:  {text}"
        p.font.size = Pt(17)
        p.font.color.rgb = WHITE
        p.space_after = Pt(14)
        if i == 0:
            p.font.bold = True
            p.font.size = Pt(22)
    add_footer(slide, "Customize funding amount and milestones before presenting")


def slide_contact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY_DARK)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), Inches(10), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(1.0))
    t.text_frame.paragraphs[0].text = "Thank You"
    t.text_frame.paragraphs[0].font.size = Pt(44)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE

    c = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(8.5), Inches(2.5))
    tf = c.text_frame
    for i, line in enumerate([
        "Panchaayat — Make it visible. Get it resolved.",
        "",
        "Demo:  [your-vercel-url].vercel.app",
        "GitHub: github.com/himsinghvi/panchayat",
        "Email:  [founder@email.com]  |  [LinkedIn]",
        "",
        "Questions?",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18 if i < 2 else 16)
        p.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        p.space_after = Pt(6)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_market(prs)
    slide_solution(prs)
    slide_how_it_works(prs)
    slide_differentiator(prs)
    slide_product(prs)
    slide_personas(prs)
    slide_b2b(prs)
    slide_pricing(prs)
    slide_revenue(prs)
    slide_competition(prs)
    slide_gtm(prs)
    slide_tech(prs)
    slide_metrics(prs)
    slide_roadmap(prs)
    slide_risks(prs)
    slide_team(prs)
    slide_ask(prs)
    slide_contact(prs)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
