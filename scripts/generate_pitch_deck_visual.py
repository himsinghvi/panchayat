"""Generate visual Panchaayat pitch deck with graphics, charts, and illustrations."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
ASSETS = Path(__file__).resolve().parent / "deck_assets"
OUTPUT = ROOT / "Panchaayat-Pitch-Deck-Visual.pptx"

PRIMARY = RGBColor(0x25, 0x63, 0xEB)
PRIMARY_DARK = RGBColor(0x1D, 0x4E, 0xD8)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
DARK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x10, 0xB9, 0x81)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def footer(slide, text="Panchaayat — Confidential  |  2026"):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(9.2), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = SLATE


def header(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.45)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = BLUE_LIGHT


def add_image(slide, filename, left, top, width, height=None):
    path = ASSETS / filename
    if not path.exists():
        return None
    if height:
        return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def icon_card(slide, x, y, w, h, title, desc, color=PRIMARY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.15), Inches(y + 0.15), Inches(0.45), Inches(0.45))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(x + 0.7), Inches(y + 0.12), Inches(w - 0.8), Inches(h - 0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.bold = True
    p1.font.size = Pt(12)
    p1.font.color.rgb = DARK
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(9)
    p2.font.color.rgb = SLATE


def bullets(slide, items, left=0.5, top=1.2, width=4.5, size=13):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)


# ── Slides ──────────────────────────────────────────────────────────

def slide_title(prs):
    s = blank_slide(prs)
    add_image(s, "deck_hero.png", 0, 0, 10, 7.5)
    # Dark overlay bottom
    overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.2), Inches(10), Inches(3.3))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    overlay.fill.transparency = 0.15
    overlay.line.fill.background()
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.15), Inches(10), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.55), Inches(4.45), Inches(8), Inches(0.8))
    t.text_frame.paragraphs[0].text = "Panchaayat"
    t.text_frame.paragraphs[0].font.size = Pt(48)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    sub = s.shapes.add_textbox(Inches(0.55), Inches(5.35), Inches(8.5), Inches(0.9))
    sub.text_frame.paragraphs[0].text = "India's public record for consumer resolutions — and case management SMEs can trust"
    sub.text_frame.paragraphs[0].font.size = Pt(18)
    sub.text_frame.paragraphs[0].font.color.rgb = BLUE_LIGHT
    tag = s.shapes.add_textbox(Inches(0.55), Inches(6.35), Inches(6), Inches(0.4))
    tag.text_frame.paragraphs[0].text = "Investor & Partnership Pitch"
    tag.text_frame.paragraphs[0].font.size = Pt(13)
    tag.text_frame.paragraphs[0].font.color.rgb = ACCENT


def slide_problem(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "The Problem", "Fragmented complaints. No accountability. No closed loop.")
    add_image(s, "deck_problem.png", 0.4, 1.2, 5.8, 3.8)
    bullets(s, [
        "Reviews, complaints & grievances live in silos",
        "Brands self-mark issues 'resolved' without proof",
        "SMEs lose queries in shared email inboxes",
        "Consumers can't see real resolution track records",
    ], left=6.4, top=1.5, width=3.2, size=12)
    footer(s)


def slide_market(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "Market Opportunity", "India-first — massive consumer economy + underserved SME support")
    add_image(s, "chart_market.png", 0.35, 1.15, 9.3, 4.5)
    callout = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.85), Inches(9), Inches(0.75))
    callout.fill.solid()
    callout.fill.fore_color.rgb = BLUE_LIGHT
    callout.line.color.rgb = PRIMARY
    callout.text_frame.paragraphs[0].text = "Wedge: SEO on '[Brand] complaints' + affordable SME SaaS vs. Zendesk/Freshdesk"
    callout.text_frame.paragraphs[0].font.size = Pt(13)
    callout.text_frame.paragraphs[0].font.color.rgb = DARK
    callout.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer(s)


def slide_solution(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "The Solution", "One platform — public reputation + private case management")
    # Hub diagram
    cx, cy = 5.0, 3.8
    hub = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.9), Inches(cy - 0.55), Inches(1.8), Inches(1.1))
    hub.fill.solid()
    hub.fill.fore_color.rgb = PRIMARY
    hub.line.fill.background()
    hub.text_frame.paragraphs[0].text = "Panchaayat\nHub"
    hub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    hub.text_frame.paragraphs[0].font.size = Pt(14)
    hub.text_frame.paragraphs[0].font.bold = True
    hub.text_frame.paragraphs[0].font.color.rgb = WHITE
    nodes = [
        (1.2, 2.0, "Public\nComplaints", ACCENT),
        (7.8, 2.0, "Brand\nProfiles", GREEN),
        (1.2, 5.2, "SME Case\nManagement", PRIMARY),
        (7.8, 5.2, "AI + API\nLayer", RGBColor(0x63, 0x66, 0xF1)),
    ]
    for nx, ny, label, col in nodes:
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(nx), Inches(ny), Inches(1.5), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = col
        box.line.fill.background()
        box.text_frame.paragraphs[0].text = label
        box.text_frame.paragraphs[0].font.size = Pt(11)
        box.text_frame.paragraphs[0].font.bold = True
        box.text_frame.paragraphs[0].font.color.rgb = WHITE
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # connector line
        line = s.shapes.add_connector(1, Inches(nx + 0.75), Inches(ny + 0.42), Inches(cx), Inches(cy))
        line.line.color.rgb = SLATE
        line.line.width = Pt(1.5)
    tag = s.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.4))
    tag.text_frame.paragraphs[0].text = "Make it visible. Get it resolved. Verified by you."
    tag.text_frame.paragraphs[0].font.size = Pt(16)
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    footer(s)


def slide_workflow(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "How It Works", "Consumer-confirmed resolution — our core differentiator")
    add_image(s, "chart_workflow.png", 0.3, 1.2, 9.4, 3.2)
    bullets(s, [
        "Guest or registered users post with evidence",
        "Community discusses; Me Too builds pressure",
        "Brands respond officially on the public record",
        "Only the original consumer can close the case",
    ], left=0.5, top=4.6, width=9, size=12)
    footer(s)


def slide_product(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "Product — MVP Live", "Full-stack platform deployed and demo-ready")
    features = [
        (0.4, 1.25, "Complaint Wizard", "AI draft & quality check", PRIMARY),
        (2.55, 1.25, "Smart Search", "NL queries + synonyms", GREEN),
        (4.7, 1.25, "Brand Dashboard", "KPIs + inbox", ACCENT),
        (6.85, 1.25, "Resolution Loop", "Propose → confirm", RGBColor(0x63, 0x66, 0xF1)),
        (0.4, 3.1, "Public Feed", "Trending + resolved", PRIMARY),
        (2.55, 3.1, "Admin + Ads", "AI targeting", GREEN),
        (4.7, 3.1, "B2B Pages", "Pricing + API docs", ACCENT),
        (6.85, 3.1, "Themes", "Dark / light mode", RGBColor(0x63, 0x66, 0xF1)),
    ]
    for x, y, title, desc, col in features:
        icon_card(s, x, y, 2.0, 1.55, title, desc, col)
    demo = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(4.9), Inches(9.2), Inches(1.5))
    demo.fill.solid()
    demo.fill.fore_color.rgb = DARK
    demo.line.fill.background()
    tf = demo.text_frame
    tf.paragraphs[0].text = "Live Demo"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ACCENT
    p2 = tf.add_paragraph()
    p2.text = "github.com/himsinghvi/panchayat  •  Login: admin / demo123  •  8 personas, 6 brands, seeded complaints"
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE
    footer(s)


def slide_b2b(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "B2B — For Business", "Turn your inbox into a case management system")
    add_image(s, "deck_b2b.png", 0.35, 1.15, 5.5, 3.5)
    bullets(s, [
        "Branded page: panchaayat.in/yourbrand",
        "2-way email-to-case sync",
        "Case ID, status, owner for every query",
        "Private by default; public escalation if unresolved",
        "CSAT/NPS + SLA dashboards",
        "14-day free trial",
    ], left=6.1, top=1.4, width=3.5, size=12)
    footer(s)


def slide_pricing(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "SaaS Pricing", "Simple plans for SMEs and MSMEs")
    plans = [
        ("Startup", "₹1,999", "/mo", ["150 cases/mo", "1 seat", "Branded page", "Email + SMS"], False),
        ("Scale-up", "₹3,999", "/mo", ["Unlimited cases", "10 seats", "Email sync", "CSAT + social"], True),
        ("Enterprise", "₹5,999", "/mo", ["Full API", "White-label", "AI triage", "Unlimited seats"], False),
    ]
    for i, (name, price, unit, feats, featured) in enumerate(plans):
        x = 0.45 + i * 3.15
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(2.85), Inches(5.2))
        card.fill.solid()
        card.fill.fore_color.rgb = BLUE_LIGHT if featured else WHITE
        card.line.color.rgb = PRIMARY if featured else RGBColor(0xE2, 0xE8, 0xF0)
        card.line.width = Pt(2.5 if featured else 1)
        if featured:
            badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.6), Inches(1.15), Inches(1.65), Inches(0.35))
            badge.fill.solid()
            badge.fill.fore_color.rgb = ACCENT
            badge.line.fill.background()
            badge.text_frame.paragraphs[0].text = "MOST POPULAR"
            badge.text_frame.paragraphs[0].font.size = Pt(8)
            badge.text_frame.paragraphs[0].font.bold = True
            badge.text_frame.paragraphs[0].font.color.rgb = WHITE
            badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tb = s.shapes.add_textbox(Inches(x + 0.2), Inches(1.55), Inches(2.45), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, (text, bold, size, color) in enumerate([
            (name, True, 18, PRIMARY),
            (f"{price}{unit}", True, 26, DARK),
            ("", False, 6, DARK),
        ] + [(f"✓  {f}", False, 11, SLATE) for f in feats]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = text
            p.font.bold = bold
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.space_after = Pt(4)
    footer(s)


def slide_revenue(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "Business Model", "Trust-first monetization")
    add_image(s, "chart_revenue.png", 0.3, 1.1, 4.8, 4.8)
    bullets(s, [
        "SaaS subscriptions — primary revenue (Startup → Enterprise)",
        "Enterprise API & white-label — banks, insurers, telcos",
        "Contextual ads — legal aid, warranty, refund tools (never bias complaints)",
        "Future: verified analytics, researcher API",
        "Explicitly avoided: pay-to-remove complaints",
    ], left=5.3, top=1.4, width=4.3, size=12)
    footer(s)


def slide_competition(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "Competitive Landscape", "No incumbent owns the full resolution loop")
    add_image(s, "chart_comparison.png", 0.25, 1.1, 9.5, 5.0)
    footer(s)


def slide_gtm(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "Go-to-Market", "Hyper-local density → national scale")
    phases = [
        ("Phase 1", "Seed Pune/Hyderabad\nwith demo brands", PRIMARY),
        ("Phase 2", "SEO: [Brand] complaints\norganic traffic", GREEN),
        ("Phase 3", "SME outbound\nD2C & appliances", ACCENT),
        ("Phase 4", "Enterprise API\nbanks & insurers", RGBColor(0x63, 0x66, 0xF1)),
    ]
    for i, (phase, desc, col) in enumerate(phases):
        x = 0.5 + i * 2.35
        arrow = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(1.5), Inches(2.1), Inches(1.8))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = col
        arrow.line.fill.background()
        tb = s.shapes.add_textbox(Inches(x + 0.15), Inches(1.7), Inches(1.8), Inches(1.4))
        tf = tb.text_frame
        tf.paragraphs[0].text = phase
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = WHITE
    bullets(s, [
        "Community growth via Me Too, trending feed, recently resolved",
        "14-day free trial for business accounts",
        "Partnerships with consumer bodies (NCH, E-Jagriti) in roadmap",
    ], left=0.5, top=3.6, width=9, size=12)
    footer(s)


def slide_metrics(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "Traction & KPIs", "Where we are today — what we measure")
    add_image(s, "chart_metrics.png", 0.3, 1.15, 9.4, 2.2)
    kpis = [
        ("Brand response\n< 7 days", "North-star"),
        ("Consumer-confirmed\nresolution %", "North-star"),
        ("Median time-to-\nresolution", "North-star"),
        ("SME trial → paid\nconversion", "GTM"),
    ]
    for i, (kpi, tag) in enumerate(kpis):
        x = 0.5 + i * 2.35
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.6), Inches(2.1), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = PRIMARY
        tb = s.shapes.add_textbox(Inches(x + 0.1), Inches(3.7), Inches(1.9), Inches(1.3))
        tf = tb.text_frame
        tf.paragraphs[0].text = kpi
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = tag
        p2.font.size = Pt(9)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER
    footer(s)


def slide_roadmap(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "Roadmap", "12–18 month product & GTM plan")
    add_image(s, "chart_roadmap.png", 0.25, 1.15, 9.5, 4.2)
    footer(s)


def slide_risks(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    header(s, "Risks & Mitigations", "Honest assessment for investors")
    risks = [
        ("Defamation", "ToS, evidence reqs, takedown process", RGBColor(0xEF, 0x44, 0x44)),
        ("Fake reviews", "Verified purchase, moderation queue", ACCENT),
        ("Cold start", "Hyper-local seeding, SEO pages", PRIMARY),
        ("Brand incentive", "Public pages rank on Google", GREEN),
        ("SME churn", "Email sync lock-in, 30-day guarantee", RGBColor(0x63, 0x66, 0xF1)),
    ]
    for i, (risk, mit, col) in enumerate(risks):
        y = 1.3 + i * 1.1
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(2.2), Inches(0.85))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        bar.text_frame.paragraphs[0].text = risk
        bar.text_frame.paragraphs[0].font.size = Pt(12)
        bar.text_frame.paragraphs[0].font.bold = True
        bar.text_frame.paragraphs[0].font.color.rgb = WHITE
        bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        mit_box = s.shapes.add_textbox(Inches(2.8), Inches(y + 0.1), Inches(6.8), Inches(0.7))
        mit_box.text_frame.paragraphs[0].text = f"→  {mit}"
        mit_box.text_frame.paragraphs[0].font.size = Pt(12)
        mit_box.text_frame.paragraphs[0].font.color.rgb = DARK
    footer(s)


def slide_ask(prs):
    s = blank_slide(prs)
    bg(s, PRIMARY_DARK)
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8), Inches(0.7))
    t.text_frame.paragraphs[0].text = "The Ask"
    t.text_frame.paragraphs[0].font.size = Pt(32)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    blocks = [
        ("₹[X] Cr Seed Round", "Customize before presenting", ACCENT),
        ("40% GTM & Sales", "30% Engineering  |  20% Ops & Legal  |  10% Buffer", PRIMARY),
        ("12-mo Milestones", "100 paying SMEs  •  10K MAU  •  Enterprise pilot", GREEN),
    ]
    for i, (title, desc, col) in enumerate(blocks):
        y = 1.5 + i * 1.65
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y), Inches(8.8), Inches(1.35))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
        card.line.color.rgb = col
        card.line.width = Pt(2)
        tb = s.shapes.add_textbox(Inches(0.85), Inches(y + 0.2), Inches(8.3), Inches(1))
        tf = tb.text_frame
        tf.paragraphs[0].text = title
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = col
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = BLUE_LIGHT
    footer(s, "Customize funding amount before investor meetings")


def slide_contact(prs):
    s = blank_slide(prs)
    add_image(s, "deck_hero.png", 0, 0, 10, 7.5)
    overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = PRIMARY_DARK
    overlay.fill.transparency = 0.25
    overlay.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(8.5), Inches(1))
    t.text_frame.paragraphs[0].text = "Thank You"
    t.text_frame.paragraphs[0].font.size = Pt(52)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = WHITE
    c = s.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.5), Inches(2.5))
    tf = c.text_frame
    for i, line in enumerate([
        "Make it visible. Get it resolved.",
        "",
        "Demo:  [your-vercel-url].vercel.app",
        "GitHub: github.com/himsinghvi/panchayat",
        "Email:  [founder@email.com]",
        "",
        "Questions?",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18 if i == 0 else 15)
        p.font.color.rgb = BLUE_LIGHT


def build():
    # Ensure charts exist
    import subprocess, sys
    subprocess.run([sys.executable, str(ASSETS / "generate_charts.py")], check=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_market(prs)
    slide_solution(prs)
    slide_workflow(prs)
    slide_product(prs)
    slide_b2b(prs)
    slide_pricing(prs)
    slide_revenue(prs)
    slide_competition(prs)
    slide_gtm(prs)
    slide_metrics(prs)
    slide_roadmap(prs)
    slide_risks(prs)
    slide_ask(prs)
    slide_contact(prs)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
