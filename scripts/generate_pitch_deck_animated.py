"""Generate an ANIMATED, graphics-first Panchaayat pitch deck.

Beyond the static visual deck this adds:
  * Native vector infographics (crisp, editable) instead of text bullets
  * Auto-playing staggered ENTRANCE animations per shape (fade / wipe / rise)
  * Slide TRANSITIONS between every slide

python-pptx has no animation API, so the <p:transition> and <p:timing>
trees are injected as raw DrawingML XML.
"""
from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parent.parent
ASSETS = Path(__file__).resolve().parent / "deck_assets"
OUTPUT = ROOT / "Panchaayat-Pitch-Deck-Animated.pptx"

# ── Palette ──────────────────────────────────────────────────────────
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
PRIMARY_DARK = RGBColor(0x1D, 0x4E, 0xD8)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
DARK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
GREEN = RGBColor(0x10, 0xB9, 0x81)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
RED = RGBColor(0xEF, 0x44, 0x44)
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)


# ══════════════════════════════════════════════════════════════════════
#  ANIMATION ENGINE  —  raw XML injection
# ══════════════════════════════════════════════════════════════════════
def _tag(name):
    return f"{{http://schemas.openxmlformats.org/presentationml/2006/main}}{name}"


def _effect_children(spid, effect, dur, cid):
    """Return the <p:set> + effect behaviour XML for one shape entrance."""
    set_xml = f"""
      <p:set>
        <p:cBhvr>
          <p:cTn id="{cid[0]}" dur="1" fill="hold">
            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
          </p:cTn>
          <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
        </p:cBhvr>
        <p:to><p:strVal val="visible"/></p:to>
      </p:set>"""
    cid[0] += 1

    filt = {"fade": "fade", "wipe": "wipe(up)"}.get(effect, "fade")
    anim_effect = f"""
      <p:animEffect transition="in" filter="{filt}">
        <p:cBhvr>
          <p:cTn id="{cid[0]}" dur="{dur}"/>
          <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
        </p:cBhvr>
      </p:animEffect>"""
    cid[0] += 1

    extra = ""
    if effect == "rise":  # subtle upward float paired with the fade
        extra = f"""
      <p:anim calcmode="lin" valueType="num" additive="base">
        <p:cBhvr>
          <p:cTn id="{cid[0]}" dur="{dur}" fill="hold"/>
          <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
          <p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst>
        </p:cBhvr>
        <p:tavLst>
          <p:tav tm="0"><p:val><p:strVal val="#ppt_y+0.08"/></p:val></p:tav>
          <p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav>
        </p:tavLst>
      </p:anim>"""
        cid[0] += 1
    return set_xml + anim_effect + extra


def _preset(effect):
    return {"fade": (10, 0), "wipe": (22, 4), "rise": (10, 0)}.get(effect, (10, 0))


def animate(slide, specs):
    """specs: list of (shape, effect, delay_ms, dur_ms).

    Builds one auto-playing main sequence; each effect starts `delay_ms`
    after the previous one finishes ("After Previous")."""
    cid = [3]  # running cTn id (1,2 reserved for tmRoot + mainSeq)
    blocks = []
    for shape, effect, delay, dur in specs:
        spid = shape.shape_id
        preset_id, sub = _preset(effect)
        outer, mid, eff = cid[0], cid[0] + 1, cid[0] + 2
        cid[0] += 3
        children = _effect_children(spid, effect, dur, cid)
        blocks.append(f"""
        <p:par>
          <p:cTn id="{outer}" fill="hold">
            <p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>
            <p:childTnLst>
              <p:par>
                <p:cTn id="{mid}" fill="hold">
                  <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="{eff}" presetID="{preset_id}" presetClass="entr" presetSubtype="{sub}" fill="hold" grpId="0" nodeType="afterEffect">
                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                        <p:childTnLst>{children}
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
              </p:par>
            </p:childTnLst>
          </p:cTn>
        </p:par>""")

    timing_xml = f"""<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>{''.join(blocks)}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""
    timing = etree.fromstring(timing_xml)
    # timing must be the last child of <p:sld>
    old = slide._element.find(_tag("timing"))
    if old is not None:
        slide._element.remove(old)
    slide._element.append(timing)


def transition(slide, kind="fade", speed="med"):
    """Add a slide transition (fade / push-left / cover-up)."""
    body = {
        "fade": "<p:fade/>",
        "push": '<p:push dir="l"/>',
        "cover": '<p:cover dir="u"/>',
        "wipe": '<p:wipe dir="l"/>',
    }.get(kind, "<p:fade/>")
    xml = (f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
           f'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" spd="{speed}">'
           f'{body}</p:transition>')
    trans = etree.fromstring(xml)
    # OOXML child order: cSld, clrMapOvr, transition, timing.
    sld = slide._element
    old = sld.find(_tag("transition"))
    if old is not None:
        sld.remove(old)
    anchor = sld.find(_tag("clrMapOvr"))
    if anchor is None:
        anchor = sld.find(_tag("cSld"))
    anchor.addnext(trans)  # place transition right after clrMapOvr (or cSld)


def cascade(shapes, effect="rise", start=200, gap=180, dur=500):
    """Convenience: build animate() specs that stagger a list of shapes."""
    specs = []
    for i, sh in enumerate(shapes):
        specs.append((sh, effect, start if i == 0 else gap, dur))
    return specs


# ══════════════════════════════════════════════════════════════════════
#  DRAWING HELPERS
# ══════════════════════════════════════════════════════════════════════
def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_text(shape, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
    return shape


def textbox(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_text(tb, text, size, color, bold, align, anchor)


def rounded(slide, x, y, w, h, fill=WHITE, line=None, line_w=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(sh, fill)
    if line is None:
        _no_line(sh)
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def oval(slide, x, y, d, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    _fill(sh, fill)
    _no_line(sh)
    sh.shadow.inherit = False
    return sh


def add_image(slide, filename, left, top, width, height=None):
    path = ASSETS / filename
    if not path.exists():
        return None
    if height:
        return slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))
    return slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def header(slide, title, subtitle=""):
    """Colored top band. Returns the band shape so it can animate first."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.05))
    _fill(bar, PRIMARY)
    _no_line(bar)
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_left = Inches(0.45)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = BLUE_LIGHT
    # accent underline
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.02), Inches(10), Inches(0.06))
    _fill(stripe, ACCENT)
    _no_line(stripe)
    stripe.shadow.inherit = False
    return bar, stripe


def footer(slide, text="Panchaayat — Confidential  •  2026"):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.08), Inches(9.2), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = SLATE


def stat_tile(slide, x, y, w, h, big, label, color):
    card = rounded(slide, x, y, w, h, WHITE, CARD_BORDER, 1.5)
    textbox(slide, x + 0.1, y + 0.18, w - 0.2, h * 0.5, big, 34, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(slide, x + 0.1, y + h * 0.62, w - 0.2, h * 0.35, label, 10.5, SLATE, False, PP_ALIGN.CENTER, MSO_ANCHOR.TOP)
    return card


def icon_tile(slide, x, y, w, h, glyph, title, desc, color):
    """Vector card: colored disc + glyph, title, description."""
    card = rounded(slide, x, y, w, h, WHITE, CARD_BORDER, 1.25)
    disc = oval(slide, x + 0.22, y + 0.22, 0.62, color)
    set_text(disc, glyph, 20, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(slide, x + 0.22, y + 0.95, w - 0.4, 0.4, title, 13, DARK, True, PP_ALIGN.LEFT, MSO_ANCHOR.TOP)
    textbox(slide, x + 0.22, y + 1.35, w - 0.4, h - 1.45, desc, 10, SLATE, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP)
    return [card, disc]


# ══════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════
def slide_title(prs):
    s = blank_slide(prs)
    add_image(s, "deck_hero.png", 0, 0, 10, 7.5)
    overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.1), Inches(10), Inches(3.4))
    _fill(overlay, DARK)
    overlay.fill.fore_color.rgb = DARK
    overlay.line.fill.background()
    overlay.fill.transparency = 0.12
    overlay.shadow.inherit = False
    stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.05), Inches(10), Inches(0.06))
    _fill(stripe, ACCENT); _no_line(stripe); stripe.shadow.inherit = False
    title = textbox(s, 0.55, 4.35, 8.5, 0.9, "Panchaayat", 50, WHITE, True)
    sub = textbox(s, 0.55, 5.3, 9, 0.9,
                  "India's public record for consumer resolutions —\nand case management SMEs can trust", 18, BLUE_LIGHT)
    tag = textbox(s, 0.55, 6.5, 7, 0.4, "Investor & Partnership Pitch  •  2026", 13, ACCENT, True)
    transition(s, "fade")
    animate(s, [(stripe, "wipe", 300, 400), (title, "rise", 200, 600),
                (sub, "fade", 250, 600), (tag, "fade", 200, 500)])


def slide_problem(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "The Problem", "Fragmented complaints. No accountability. No closed loop.")
    tiles = [
        ("🧩", "Fragmented", "Reviews, complaints & grievances scattered across silos", RED),
        ("🚫", "No Proof", "Brands self-mark issues 'resolved' with no verification", ACCENT),
        ("📥", "Lost Queries", "SMEs bury customer queries in shared email inboxes", INDIGO),
        ("🔍", "No Track Record", "Consumers can't see a brand's real resolution history", PRIMARY),
    ]
    shapes = []
    for i, (g, t, d, c) in enumerate(tiles):
        x = 0.45 + i * 2.4
        shapes += icon_tile(s, x, 1.55, 2.2, 2.6, g, t, d, c)
    quote = rounded(s, 0.45, 4.55, 9.1, 1.9, DARK)
    textbox(s, 0.8, 4.8, 8.5, 0.5, "\"I complained. They said 'resolved'. Nothing changed.\"", 18, WHITE, True)
    textbox(s, 0.8, 5.55, 8.5, 0.7,
            "Existing options (Google Reviews, NCH, MouthShut) never close the loop with the actual affected consumer.",
            13, BLUE_LIGHT)
    footer(s)
    transition(s, "push")
    specs = [(bar, "wipe", 100, 400), (stripe, "fade", 100, 300)]
    specs += cascade(shapes, "rise", start=200, gap=140, dur=450)
    specs.append((quote, "fade", 300, 600))
    animate(s, specs)


def slide_market(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "Market Opportunity", "India-first: huge consumer economy + underserved SME support")
    img = add_image(s, "chart_market.png", 0.5, 1.35, 9.0)
    callout = rounded(s, 0.5, 5.95, 9.0, 0.75, BLUE_LIGHT, PRIMARY, 1.5)
    set_text(callout, "Wedge: rank on '[Brand] complaints' SEO  +  affordable SME SaaS vs. Zendesk / Freshdesk",
             13, DARK, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    footer(s)
    transition(s, "push")
    animate(s, [(bar, "wipe", 100, 400), (img, "rise", 250, 700), (callout, "fade", 300, 600)])


def slide_solution(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "The Solution", "One platform — public reputation + private case management")
    cx, cy = 4.6, 3.7
    # connectors first (so they sit behind nodes) — but animate hub last
    nodes = [
        (1.0, 2.0, "Public\nComplaints", ACCENT),
        (7.4, 2.0, "Brand\nProfiles", GREEN),
        (1.0, 5.0, "SME Case\nManagement", PRIMARY),
        (7.4, 5.0, "AI + API\nLayer", INDIGO),
    ]
    conns = []
    for nx, ny, label, col in nodes:
        line = s.shapes.add_connector(1, Inches(nx + 0.75), Inches(ny + 0.42), Inches(cx + 0.9), Inches(cy + 0.55))
        line.line.color.rgb = SLATE
        line.line.width = Pt(1.75)
        line.line.dash_style = None
        conns.append(line)
    node_shapes = []
    for nx, ny, label, col in nodes:
        box = rounded(s, nx, ny, 1.6, 0.9, col)
        set_text(box, label, 12, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        node_shapes.append(box)
    hub = oval(s, cx, cy - 0.05, 1.85, PRIMARY_DARK)
    set_text(hub, "Panchaayat\nHub", 15, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    tag = textbox(s, 0.5, 6.55, 9, 0.5, "Make it visible.  Get it resolved.  Verified by you.",
                  17, PRIMARY, True, PP_ALIGN.CENTER)
    footer(s)
    transition(s, "fade")
    specs = [(bar, "wipe", 100, 400), (hub, "rise", 250, 500)]
    specs += cascade(node_shapes, "rise", start=150, gap=160, dur=420)
    specs += [(c, "fade", 80, 300) for c in conns]
    specs.append((tag, "fade", 250, 500))
    animate(s, specs)


def slide_workflow(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "How It Works", "Consumer-confirmed resolution — our core differentiator")
    steps = [
        ("1", "Share", "post + evidence", PRIMARY),
        ("2", "Discuss", "community + Me Too", INDIGO),
        ("3", "Respond", "brand replies", RGBColor(0x0E, 0xA5, 0xE9)),
        ("4", "Propose", "refund / fix", ACCENT),
        ("5", "Confirm", "only YOU close it", GREEN),
    ]
    circles, labels, arrows = [], [], []
    y = 2.5
    for i, (n, t, d, c) in enumerate(steps):
        x = 0.7 + i * 1.85
        circ = oval(s, x, y, 1.05, c)
        set_text(circ, n, 26, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        circles.append(circ)
        lbl = textbox(s, x - 0.25, y + 1.15, 1.55, 0.85, f"{t}\n{d}", 12, DARK, True, PP_ALIGN.CENTER, MSO_ANCHOR.TOP)
        # make description slate
        lbl.text_frame.paragraphs[1].font.bold = False
        lbl.text_frame.paragraphs[1].font.size = Pt(9.5)
        lbl.text_frame.paragraphs[1].font.color.rgb = SLATE
        labels.append(lbl)
        if i < 4:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.12), Inches(y + 0.38), Inches(0.6), Inches(0.28))
            _fill(arr, SLATE); _no_line(arr); arr.shadow.inherit = False
            arrows.append(arr)
    band = rounded(s, 0.5, 5.15, 9.0, 1.4, DARK)
    textbox(s, 0.85, 5.35, 8.4, 0.5, "Brands cannot unilaterally close a complaint.", 17, ACCENT, True)
    textbox(s, 0.85, 5.95, 8.4, 0.5, "The full, immutable resolution history stays public forever — that is the trust moat.",
            12.5, BLUE_LIGHT)
    footer(s)
    transition(s, "push")
    specs = [(bar, "wipe", 100, 400)]
    for i in range(5):
        specs.append((circles[i], "rise", 150 if i else 200, 350))
        specs.append((labels[i], "fade", 60, 300))
        if i < 4:
            specs.append((arrows[i], "wipe", 60, 250))
    specs.append((band, "fade", 300, 600))
    animate(s, specs)


def slide_product(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "Product — MVP is Live", "Full-stack platform, deployed and demo-ready")
    feats = [
        ("✍", "Complaint Wizard", "AI draft + quality check", PRIMARY),
        ("🔎", "Smart Search", "NL queries + synonyms", GREEN),
        ("📊", "Brand Dashboard", "KPIs + complaint inbox", ACCENT),
        ("🔁", "Resolution Loop", "propose → confirm", INDIGO),
        ("📰", "Public Feed", "trending + resolved", PRIMARY),
        ("🎯", "Contextual Ads", "AI targeting, no bias", GREEN),
        ("🏷", "B2B Pages", "pricing + API docs", ACCENT),
        ("🤖", "AI Layer", "Azure GPT + fallbacks", INDIGO),
    ]
    tiles = []
    for i, (g, t, d, c) in enumerate(feats):
        col, row = i % 4, i // 4
        x = 0.45 + col * 2.4
        yy = 1.35 + row * 1.85
        tiles += icon_tile(s, x, yy, 2.2, 1.7, g, t, d, c)
    demo = rounded(s, 0.45, 5.15, 9.1, 1.35, DARK)
    textbox(s, 0.75, 5.32, 8.5, 0.4, "Live Demo", 14, ACCENT, True)
    textbox(s, 0.75, 5.78, 8.7, 0.65,
            "github.com/himsinghvi/panchayat   •   login  admin / demo123   •   8 personas · 6 brands · seeded cases",
            12, WHITE)
    footer(s)
    transition(s, "push")
    specs = [(bar, "wipe", 100, 400)]
    specs += cascade(tiles, "rise", start=150, gap=90, dur=380)
    specs.append((demo, "fade", 250, 600))
    animate(s, specs)


def slide_b2b(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "For Business", "Turn a chaotic inbox into a case-management system")
    img = add_image(s, "deck_b2b.png", 0.4, 1.4, 4.9)
    rows = [
        ("🌐", "Branded page  panchaayat.in/yourbrand", PRIMARY),
        ("✉", "2-way email → case sync", GREEN),
        ("🆔", "Case ID, status & owner on every query", ACCENT),
        ("🔒", "Private by default; public escalation if ignored", INDIGO),
        ("📈", "CSAT / NPS + SLA dashboards", PRIMARY),
        ("🎁", "14-day free trial, cancel anytime", GREEN),
    ]
    row_shapes = []
    for i, (g, t, c) in enumerate(rows):
        y = 1.55 + i * 0.83
        chip = rounded(s, 5.6, y, 0.62, 0.62, c)
        set_text(chip, g, 16, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        lbl = textbox(s, 6.35, y + 0.03, 3.4, 0.6, t, 12.5, DARK, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        row_shapes += [chip, lbl]
    footer(s)
    transition(s, "push")
    specs = [(bar, "wipe", 100, 400), (img, "fade", 200, 600)]
    specs += cascade(row_shapes, "rise", start=150, gap=110, dur=350)
    animate(s, specs)


def slide_pricing(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "SaaS Pricing", "Simple, affordable plans for SMEs & MSMEs")
    plans = [
        ("Startup", "₹1,999", "/mo", ["150 cases / mo", "1 seat", "Branded page", "Email + SMS"], False),
        ("Scale-up", "₹3,999", "/mo", ["Unlimited cases", "10 seats", "Email → case sync", "CSAT + social"], True),
        ("Enterprise", "₹5,999", "/mo", ["Full REST API", "White-label", "AI triage", "Unlimited seats"], False),
    ]
    cards = []
    for i, (name, price, unit, feats, featured) in enumerate(plans):
        x = 0.55 + i * 3.1
        card = rounded(s, x, 1.55, 2.8, 4.7, BLUE_LIGHT if featured else WHITE,
                       PRIMARY if featured else CARD_BORDER, 2.5 if featured else 1.25)
        if featured:
            badge = rounded(s, x + 0.65, 1.38, 1.5, 0.36, ACCENT)
            set_text(badge, "MOST POPULAR", 8.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        textbox(s, x + 0.2, 1.75, 2.4, 0.45, name, 18, PRIMARY, True, PP_ALIGN.CENTER)
        textbox(s, x + 0.2, 2.3, 2.4, 0.6, price + unit, 26, DARK, True, PP_ALIGN.CENTER)
        ftxt = "\n".join("✓  " + f for f in feats)
        textbox(s, x + 0.3, 3.15, 2.2, 2.9, ftxt, 12, SLATE, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP)
        cards.append(card)
    note = textbox(s, 0.5, 6.45, 9, 0.4, "All plans: no setup fee  •  30-day money-back  •  GST invoice",
                   11, SLATE, False, PP_ALIGN.CENTER)
    footer(s)
    transition(s, "push")
    specs = [(bar, "wipe", 100, 400)]
    specs += cascade(cards, "rise", start=200, gap=250, dur=500)
    specs.append((note, "fade", 200, 400))
    animate(s, specs)


def slide_revenue(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "Business Model", "Trust-first monetization")
    img = add_image(s, "chart_revenue.png", 0.35, 1.4, 4.7)
    rows = [
        ("SaaS subscriptions", "primary revenue — Startup → Enterprise", PRIMARY),
        ("Enterprise API & white-label", "banks, insurers, telcos", ACCENT),
        ("Contextual ads", "legal aid, warranty, refund tools — never bias complaints", GREEN),
        ("Verified analytics (future)", "trends, benchmarks, researcher API", INDIGO),
    ]
    row_shapes = []
    for i, (t, d, c) in enumerate(rows):
        y = 1.6 + i * 1.05
        chip = rounded(s, 5.35, y, 0.14, 0.85, c)
        lbl = textbox(s, 5.65, y - 0.02, 4.0, 0.5, t, 13.5, DARK, True, PP_ALIGN.LEFT, MSO_ANCHOR.TOP)
        sub = textbox(s, 5.65, y + 0.42, 4.0, 0.5, d, 10.5, SLATE, False, PP_ALIGN.LEFT, MSO_ANCHOR.TOP)
        row_shapes += [chip, lbl, sub]
    guard = rounded(s, 5.35, 5.95, 4.3, 0.6, DARK)
    set_text(guard, "✗  We never let brands pay to remove complaints", 11.5, RED, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    guard.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFC, 0xA5, 0xA5)
    footer(s)
    transition(s, "push")
    specs = [(bar, "wipe", 100, 400), (img, "rise", 200, 600)]
    specs += cascade(row_shapes, "fade", start=150, gap=90, dur=320)
    specs.append((guard, "rise", 250, 450))
    animate(s, specs)


def slide_competition(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "Competitive Landscape", "No incumbent owns the full resolution loop")
    img = add_image(s, "chart_comparison.png", 0.35, 1.35, 9.3)
    footer(s)
    transition(s, "fade")
    animate(s, [(bar, "wipe", 100, 400), (img, "rise", 250, 700)])


def slide_gtm(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "Go-to-Market", "Hyper-local density → national scale")
    phases = [
        ("Phase 1", "Seed Pune / Hyderabad\nwith demo brands", PRIMARY),
        ("Phase 2", "SEO: '[Brand] complaints'\norganic traffic", GREEN),
        ("Phase 3", "SME outbound —\nD2C & appliances", ACCENT),
        ("Phase 4", "Enterprise API —\nbanks & insurers", INDIGO),
    ]
    chevs = []
    for i, (p, d, c) in enumerate(phases):
        x = 0.5 + i * 2.35
        ch = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(1.55), Inches(2.25), Inches(1.9))
        _fill(ch, c); _no_line(ch); ch.shadow.inherit = False
        set_text(ch, f"{p}\n{d}", 11.5, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        ch.text_frame.paragraphs[0].font.size = Pt(13)
        chevs.append(ch)
    rows = [
        ("🔥", "Viral growth via Me Too, trending & recently-resolved feeds"),
        ("🎁", "14-day free trial converts SMEs; email-sync lock-in reduces churn"),
        ("🏛", "Roadmap: partner with NCH / E-Jagriti for escalations"),
    ]
    row_shapes = []
    for i, (g, t) in enumerate(rows):
        y = 3.9 + i * 0.9
        chip = rounded(s, 0.5, y, 0.62, 0.62, PRIMARY)
        set_text(chip, g, 15, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        lbl = textbox(s, 1.25, y + 0.02, 8.3, 0.6, t, 13, DARK, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        row_shapes += [chip, lbl]
    footer(s)
    transition(s, "push")
    specs = [(bar, "wipe", 100, 400)]
    specs += cascade(chevs, "wipe", start=200, gap=220, dur=450)
    specs += cascade(row_shapes, "fade", start=150, gap=110, dur=320)
    animate(s, specs)


def slide_metrics(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "Traction & KPIs", "Where we are today — and what we measure")
    tiles_data = [("MVP", "Live & deployed", PRIMARY), ("8", "Demo personas", ACCENT),
                  ("6", "Seeded brands", GREEN), ("3", "SaaS tiers", INDIGO)]
    tiles = []
    for i, (b, l, c) in enumerate(tiles_data):
        x = 0.5 + i * 2.35
        tiles.append(stat_tile(s, x, 1.4, 2.1, 1.55, b, l, c))
    textbox(s, 0.5, 3.25, 9, 0.4, "North-star metrics we will report to investors", 13, DARK, True)
    kpis = [
        ("⏱", "Brand response < 7 days", ACCENT),
        ("✅", "Consumer-confirmed resolution %", GREEN),
        ("📉", "Median time-to-resolution", PRIMARY),
        ("💳", "SME trial → paid conversion", INDIGO),
    ]
    kpi_shapes = []
    for i, (g, t, c) in enumerate(kpis):
        y = 3.8 + i * 0.75
        chip = rounded(s, 0.5, y, 0.55, 0.55, c)
        set_text(chip, g, 14, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        lbl = textbox(s, 1.2, y + 0.02, 8.3, 0.5, t, 13, DARK, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        kpi_shapes += [chip, lbl]
    footer(s)
    transition(s, "push")
    specs = [(bar, "wipe", 100, 400)]
    specs += cascade(tiles, "rise", start=200, gap=160, dur=450)
    specs += cascade(kpi_shapes, "fade", start=150, gap=90, dur=300)
    animate(s, specs)


def slide_roadmap(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "Roadmap", "12–18 month product & GTM plan")
    img = add_image(s, "chart_roadmap.png", 0.35, 1.5, 9.3)
    footer(s)
    transition(s, "fade")
    animate(s, [(bar, "wipe", 100, 400), (img, "rise", 250, 700)])


def slide_risks(prs):
    s = blank_slide(prs)
    bg(s, LIGHT_BG)
    bar, stripe = header(s, "Risks & Mitigations", "Honest assessment for investors")
    risks = [
        ("Defamation", "ToS, evidence requirement, takedown & appeal process", RED),
        ("Fake reviews", "Verified-purchase signal + moderation queue", ACCENT),
        ("Cold start", "Hyper-local seeding + SEO complaint pages", PRIMARY),
        ("Brand apathy", "Public pages rank on Google for '[brand] complaints'", GREEN),
        ("SME churn", "Email-sync lock-in + 30-day money-back guarantee", INDIGO),
    ]
    shapes = []
    for i, (risk, mit, c) in enumerate(risks):
        y = 1.4 + i * 1.05
        bar_sh = rounded(s, 0.45, y, 2.3, 0.85, c)
        set_text(bar_sh, risk, 13, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.85), Inches(y + 0.28), Inches(0.5), Inches(0.28))
        _fill(arrow, SLATE); _no_line(arrow); arrow.shadow.inherit = False
        mit_card = rounded(s, 3.5, y, 6.05, 0.85, WHITE, CARD_BORDER, 1.25)
        set_text(mit_card, mit, 12.5, DARK, False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
        mit_card.text_frame.margin_left = Inches(0.2)
        shapes += [bar_sh, arrow, mit_card]
    footer(s)
    transition(s, "push")
    specs = [(bar, "wipe", 100, 400)]
    specs += cascade(shapes, "fade", start=200, gap=100, dur=350)
    animate(s, specs)


def slide_ask(prs):
    s = blank_slide(prs)
    bg(s, PRIMARY_DARK)
    title = textbox(s, 0.6, 0.5, 8, 0.8, "The Ask", 34, WHITE, True)
    blocks = [
        ("₹[X] Cr Seed Round", "Customize this figure before presenting", ACCENT),
        ("Use of funds", "40% GTM & Sales  •  30% Engineering  •  20% Ops & Legal  •  10% Buffer", PRIMARY),
        ("12-month milestones", "100 paying SMEs   •   10K MAU   •   1 enterprise pilot", GREEN),
    ]
    cards = []
    for i, (t, d, c) in enumerate(blocks):
        y = 1.55 + i * 1.6
        card = rounded(s, 0.6, y, 8.8, 1.35, RGBColor(0x1E, 0x29, 0x3B), c, 2.0)
        textbox(s, 0.9, y + 0.2, 8.3, 0.5, t, 20, c, True)
        textbox(s, 0.9, y + 0.72, 8.3, 0.5, d, 13, BLUE_LIGHT)
        cards.append(card)
    footer(s, "Customize funding amount & milestones before investor meetings")
    transition(s, "cover")
    specs = [(title, "rise", 200, 500)]
    specs += cascade(cards, "rise", start=200, gap=280, dur=550)
    animate(s, specs)


def slide_contact(prs):
    s = blank_slide(prs)
    add_image(s, "deck_hero.png", 0, 0, 10, 7.5)
    overlay = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    _fill(overlay, PRIMARY_DARK); _no_line(overlay); overlay.fill.transparency = 0.22
    overlay.shadow.inherit = False
    ty = textbox(s, 0.7, 2.1, 8.6, 1.0, "Thank You", 52, WHITE, True)
    tag = textbox(s, 0.7, 3.25, 8.6, 0.6, "Make it visible.  Get it resolved.", 20, ACCENT, True)
    contact = textbox(s, 0.7, 4.2, 8.6, 2.0,
                      "Demo:   [your-vercel-url].vercel.app\n"
                      "GitHub:  github.com/himsinghvi/panchayat\n"
                      "Email:   [founder@email.com]\n\nQuestions?",
                      16, BLUE_LIGHT)
    transition(s, "fade")
    animate(s, [(ty, "rise", 250, 600), (tag, "fade", 250, 500), (contact, "fade", 300, 600)])


def build():
    import subprocess, sys
    subprocess.run([sys.executable, str(ASSETS / "generate_charts.py")], check=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for fn in (slide_title, slide_problem, slide_market, slide_solution, slide_workflow,
               slide_product, slide_b2b, slide_pricing, slide_revenue, slide_competition,
               slide_gtm, slide_metrics, slide_roadmap, slide_risks, slide_ask, slide_contact):
        fn(prs)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT} ({len(prs.slides)} slides, animated)")


if __name__ == "__main__":
    build()
