"""Panchaayat pitch deck — a fresh, self-contained brand & template.

Design language (built from scratch here, reuses no prior assets):
  * 16:9 widescreen, warm-paper / deep-ink surfaces
  * Palette: teal (resolution) + terracotta coral (people) + warm gold
  * Type: Georgia (serif display headlines) + Segoe UI (body)
  * Custom "council ring" emblem as recurring chrome
  * 100% native vector graphics — no images, no external charts
  * Staggered auto-play entrance animations + slide transitions
    (injected as raw DrawingML XML — python-pptx has no animation API)
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "Panchaayat-Deck.pptx"

# ── Brand palette ────────────────────────────────────────────────────
INK      = RGBColor(0x12, 0x17, 0x2B)   # deep charcoal-navy (dark bg)
INK_SOFT = RGBColor(0x1E, 0x25, 0x3F)   # raised dark surface
PAPER    = RGBColor(0xFA, 0xF6, 0xF0)   # warm off-white (light bg)
PAPER_2  = RGBColor(0xF1, 0xEA, 0xDF)   # warm card tint
TEAL     = RGBColor(0x0D, 0x94, 0x88)   # primary brand
TEAL_DK  = RGBColor(0x0F, 0x76, 0x6E)
TEAL_LT  = RGBColor(0xCC, 0xE9, 0xE5)
CORAL    = RGBColor(0xF2, 0x6A, 0x4B)   # accent — people / warmth
CORAL_LT = RGBColor(0xF9, 0xD8, 0xCF)
GOLD     = RGBColor(0xE0, 0xA8, 0x2E)   # secondary accent
PLUM     = RGBColor(0x6D, 0x5D, 0xD3)   # tertiary
INKTEXT  = RGBColor(0x1B, 0x21, 0x33)   # body text on paper
MUTED    = RGBColor(0x7A, 0x73, 0x66)   # warm gray text
LINE     = RGBColor(0xE3, 0xDA, 0xCC)   # hairline on paper
CREAM    = RGBColor(0xF3, 0xEE, 0xE6)   # text on dark
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED      = RGBColor(0xDC, 0x3E, 0x3E)

HEAD = "Georgia"      # serif display
BODY = "Segoe UI"     # sans body

EMU_IN = 914400


# ══════════════════════════════════════════════════════════════════════
#  ANIMATION ENGINE  (raw XML injection)
# ══════════════════════════════════════════════════════════════════════
def _tag(name):
    return f"{{http://schemas.openxmlformats.org/presentationml/2006/main}}{name}"


def _effect_children(spid, effect, dur, cid):
    set_xml = f"""
      <p:set>
        <p:cBhvr>
          <p:cTn id="{cid[0]}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
          <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
          <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
        </p:cBhvr>
        <p:to><p:strVal val="visible"/></p:to>
      </p:set>"""
    cid[0] += 1
    filt = {"fade": "fade", "wipe": "wipe(up)"}.get(effect, "fade")
    anim_effect = f"""
      <p:animEffect transition="in" filter="{filt}">
        <p:cBhvr><p:cTn id="{cid[0]}" dur="{dur}"/><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>
      </p:animEffect>"""
    cid[0] += 1
    extra = ""
    if effect == "rise":
        extra = f"""
      <p:anim calcmode="lin" valueType="num" additive="base">
        <p:cBhvr><p:cTn id="{cid[0]}" dur="{dur}" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
          <p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst></p:cBhvr>
        <p:tavLst>
          <p:tav tm="0"><p:val><p:strVal val="#ppt_y+0.07"/></p:val></p:tav>
          <p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav>
        </p:tavLst>
      </p:anim>"""
        cid[0] += 1
    return set_xml + anim_effect + extra


def _preset(effect):
    return {"fade": (10, 0), "wipe": (22, 4), "rise": (10, 0)}.get(effect, (10, 0))


def animate(slide, specs):
    """specs: list of (shape, effect, delay_ms, dur_ms) — auto-play cascade."""
    cid = [3]
    blocks = []
    for shape, effect, delay, dur in specs:
        spid = shape.shape_id
        pid, sub = _preset(effect)
        outer, mid, eff = cid[0], cid[0] + 1, cid[0] + 2
        cid[0] += 3
        children = _effect_children(spid, effect, dur, cid)
        blocks.append(f"""
        <p:par><p:cTn id="{outer}" fill="hold"><p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>
          <p:childTnLst><p:par><p:cTn id="{mid}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst><p:par>
              <p:cTn id="{eff}" presetID="{pid}" presetClass="entr" presetSubtype="{sub}" fill="hold" grpId="0" nodeType="afterEffect">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst>{children}
                </p:childTnLst></p:cTn>
            </p:par></p:childTnLst></p:cTn></p:par></p:childTnLst>
        </p:cTn></p:par>""")
    timing_xml = f"""<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
    <p:seq concurrent="1" nextAc="seek">
      <p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>{''.join(blocks)}
      </p:childTnLst></p:cTn>
      <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
      <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
    </p:seq>
  </p:childTnLst></p:cTn></p:par></p:tnLst>
</p:timing>"""
    timing = etree.fromstring(timing_xml)
    old = slide._element.find(_tag("timing"))
    if old is not None:
        slide._element.remove(old)
    slide._element.append(timing)


def transition(slide, kind="fade", speed="med"):
    body = {"fade": "<p:fade/>", "push": '<p:push dir="l"/>',
            "cover": '<p:cover dir="u"/>', "wipe": '<p:wipe dir="l"/>'}.get(kind, "<p:fade/>")
    xml = (f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
           f'spd="{speed}">{body}</p:transition>')
    trans = etree.fromstring(xml)
    sld = slide._element
    old = sld.find(_tag("transition"))
    if old is not None:
        sld.remove(old)
    anchor = sld.find(_tag("clrMapOvr"))
    if anchor is None:
        anchor = sld.find(_tag("cSld"))
    anchor.addnext(trans)   # order: cSld, clrMapOvr, transition, timing


def cascade(shapes, effect="rise", start=180, gap=140, dur=440):
    return [(sh, effect, start if i == 0 else gap, dur) for i, sh in enumerate(shapes)]


# ══════════════════════════════════════════════════════════════════════
#  DRAWING PRIMITIVES
# ══════════════════════════════════════════════════════════════════════
def slide_of(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _run_font(run, name, size, color, bold, tracking):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if tracking is not None:
        run._r.get_or_add_rPr().set("spc", str(int(tracking * 100)))


def set_text(shape, text, size, color, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE, font=BODY, tracking=None, leading=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if leading:
            p.line_spacing = leading
        r = p.add_run()
        r.text = ln
        _run_font(r, font, size, color, bold, tracking)
    return shape


def textbox(slide, x, y, w, h, text, size, color, **kw):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    return set_text(tb, text, size, color, **kw)


def _shape(slide, kind, x, y, w, h, fill=None, line=None, line_w=1.0):
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    return sh


def rrect(slide, x, y, w, h, fill=WHITE, line=None, line_w=1.0, radius=0.10):
    sh = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line, line_w)
    try:  # set corner radius
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def rect(slide, x, y, w, h, fill, line=None, line_w=1.0):
    return _shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line, line_w)


def disc(slide, cx, cy, r, fill=None, line=None, line_w=1.0):
    return _shape(slide, MSO_SHAPE.OVAL, cx - r, cy - r, 2 * r, 2 * r, fill, line, line_w)


def connector(slide, x1, y1, x2, y2, color=MUTED, w=1.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    return c


# ── Custom emblem: the "council ring" ────────────────────────────────
def emblem(slide, cx, cy, r, ring=TEAL, center=CORAL, dots=TEAL, on_dark=False):
    parts = []
    outer = disc(slide, cx, cy, r, fill=None, line=ring, line_w=max(1.6, r * 8))
    parts.append(outer)
    # 4 assembly dots around the inner ring (N, E, S, W)
    dr = r * 0.62
    ddot = r * 0.14
    for ang in (0, 1, 2, 3):
        import math
        # avoid Math.random / time; deterministic positions
        a = [(-90), 0, 90, 180][ang] * 3.14159265 / 180.0
        dx = cx + dr * math.cos(a)
        dy = cy + dr * math.sin(a)
        parts.append(disc(slide, dx, dy, ddot, fill=dots))
    parts.append(disc(slide, cx, cy, r * 0.22, fill=center))
    return parts


def wordmark(slide, x, y, color_main=INKTEXT, color_dot=CORAL, size=15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.2), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r1 = p.add_run(); r1.text = "PANCHAAYAT"
    _run_font(r1, HEAD, size, color_main, True, 3)
    r2 = p.add_run(); r2.text = " ●"
    _run_font(r2, BODY, size, color_dot, True, None)
    return tb


# ── Page chrome ──────────────────────────────────────────────────────
def frame(slide, kicker, title, accent=TEAL):
    """Editorial header on a paper slide. Returns (chrome_shapes)."""
    bg(slide, PAPER)
    em = emblem(slide, 0.95, 0.78, 0.20)
    k = textbox(slide, 1.35, 0.5, 8, 0.35, kicker.upper(), 12.5, accent, bold=True,
                font=BODY, tracking=3, anchor=MSO_ANCHOR.MIDDLE)
    t = textbox(slide, 0.9, 0.92, 11.4, 0.95, title, 31, INKTEXT, bold=True, font=HEAD)
    rule = rect(slide, 0.95, 1.82, 0.85, 0.055, accent)
    chrome = em + [k, t, rule]
    return chrome


def pagechrome(slide, n, dark=False):
    cmain = CREAM if dark else INKTEXT
    cdot = CORAL
    wordmark(slide, 0.55, 7.06, color_main=(CREAM if dark else MUTED), color_dot=cdot, size=12)
    tb = textbox(slide, 12.15, 7.02, 1.0, 0.4, f"{n:02d} / 16", 10, (CREAM if dark else MUTED),
                 align=PP_ALIGN.RIGHT, font=BODY, tracking=1)


def pill(slide, x, y, w, h, text, fill, tcolor, size=10.5, bold=True):
    p = rrect(slide, x, y, w, h, fill=fill, radius=0.5)
    set_text(p, text, size, tcolor, bold=bold, align=PP_ALIGN.CENTER, font=BODY, tracking=1)
    return p


def icon_disc(slide, cx, cy, r, glyph, fill):
    d = disc(slide, cx, cy, r, fill=fill)
    set_text(d, glyph, int(r * 44), WHITE, bold=True, align=PP_ALIGN.CENTER)
    return d


# ══════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════
def s_title(prs):
    s = slide_of(prs)
    bg(s, INK)
    # faint council-ring motif, large, bottom-right
    disc(s, 11.7, 6.6, 2.6, fill=None, line=INK_SOFT, line_w=2.0)
    disc(s, 11.7, 6.6, 1.9, fill=None, line=INK_SOFT, line_w=2.0)
    # accent dots scattered (deterministic)
    for (dx, dy, c, r) in [(1.4, 1.2, TEAL, 0.06), (12.4, 1.5, CORAL, 0.07),
                            (0.9, 5.9, GOLD, 0.05), (11.2, 2.0, PLUM, 0.05)]:
        disc(s, dx, dy, r, fill=c)
    em = emblem(s, 1.15, 1.15, 0.42, ring=TEAL, center=CORAL, dots=CREAM)
    kick = textbox(s, 0.75, 2.75, 10, 0.4, "CONSUMER VOICE  ·  RESOLUTION  ·  TRUST",
                   13, TEAL, bold=True, font=BODY, tracking=3)
    name = textbox(s, 0.7, 3.15, 11.5, 1.4, "Panchaayat", 78, CREAM, bold=True, font=HEAD)
    dot = textbox(s, 6.55, 3.35, 1, 1.2, "●", 30, CORAL, font=BODY)
    sub = textbox(s, 0.78, 4.75, 11.2, 1.1,
                  "India's public record for consumer resolutions —\nand the case-management system SMEs can trust.",
                  21, CREAM, font=BODY, leading=1.15)
    line = rect(s, 0.82, 6.35, 1.1, 0.06, CORAL)
    tag = textbox(s, 0.82, 6.55, 11, 0.5, "Investor & Partnership Pitch   ·   2026", 14, MUTED,
                  font=BODY, tracking=1)
    transition(s, "fade")
    animate(s, cascade(em, "fade", 150, 60, 300) +
            [(kick, "fade", 150, 450), (name, "rise", 120, 700), (dot, "fade", 120, 400),
             (sub, "fade", 220, 600), (line, "wipe", 150, 400), (tag, "fade", 150, 500)])


def s_problem(prs):
    s = slide_of(prs)
    chrome = frame(s, "The Problem", "Complaints go out. Resolutions never come back.")
    tiles = [
        ("🧩", "Fragmented", "Reviews, complaints & grievances scattered across silos.", CORAL),
        ("🚫", "No proof", "Brands self-declare “resolved” with nothing to verify it.", GOLD),
        ("📥", "Lost queries", "SMEs bury customer issues in a shared email inbox.", PLUM),
        ("🕳", "No track record", "Consumers can’t see a brand’s real resolution history.", TEAL),
    ]
    anim = []
    for i, (g, t, d, c) in enumerate(tiles):
        x = 0.9 + i * 2.98
        card = rrect(s, x, 2.15, 2.72, 2.5, fill=WHITE, line=LINE, line_w=1.25, radius=0.09)
        icon_disc(s, x + 0.62, 2.78, 0.36, g, c)
        textbox(s, x + 0.28, 3.2, 2.2, 0.4, t, 15, INKTEXT, bold=True, font=HEAD)
        textbox(s, x + 0.28, 3.62, 2.24, 0.95, d, 11, MUTED, font=BODY, leading=1.1)
        anim.append(card)
    band = rrect(s, 0.9, 5.0, 11.52, 1.55, fill=INK, radius=0.06)
    textbox(s, 1.3, 5.18, 10.8, 0.55, "“I complained. They marked it resolved. Nothing actually changed.”",
            19, CREAM, bold=True, font=HEAD)
    textbox(s, 1.3, 5.82, 10.8, 0.6,
            "Google Reviews, NCH and MouthShut all miss the one thing that matters — the affected consumer confirming it’s truly fixed.",
            12.5, TEAL_LT, font=BODY)
    pagechrome(s, 2)
    transition(s, "push")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + cascade(anim, "rise", 200, 150, 450) +
            [(band, "fade", 250, 600)])


def s_solution(prs):
    s = slide_of(prs)
    chrome = frame(s, "The Solution", "One platform. Public reputation + private resolution.")
    cx, cy = 6.66, 4.35
    nodes = [
        (2.35, 2.95, "Public\ncomplaints", CORAL),
        (10.95, 2.95, "Brand\nprofiles", GOLD),
        (2.35, 5.75, "SME case\nmanagement", TEAL),
        (10.95, 5.75, "AI + API\nlayer", PLUM),
    ]
    conns = [connector(s, nx, ny, cx, cy, color=LINE, w=1.75) for nx, ny, _, _ in nodes]
    node_sh = []
    for nx, ny, label, c in nodes:
        b = rrect(s, nx - 1.05, ny - 0.5, 2.1, 1.0, fill=WHITE, line=c, line_w=2.0, radius=0.14)
        set_text(b, label, 12.5, INKTEXT, bold=True, align=PP_ALIGN.CENTER, font=BODY, leading=1.0)
        node_sh.append(b)
    hub = disc(s, cx, cy, 1.02, fill=TEAL_DK)
    set_text(hub, "Panchaayat", 14, WHITE, bold=True, align=PP_ALIGN.CENTER, font=HEAD)
    ring = disc(s, cx, cy, 1.18, fill=None, line=CORAL, line_w=2.0)
    tag = textbox(s, 0.9, 6.55, 11.5, 0.5, "Make it visible.   Get it resolved.   Verified by you.",
                  18, TEAL_DK, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
    pagechrome(s, 3)
    transition(s, "fade")
    animate(s, cascade(chrome, "fade", 100, 50, 300) +
            [(ring, "fade", 150, 400), (hub, "rise", 100, 500)] +
            cascade(node_sh, "rise", 150, 140, 420) +
            [(c, "fade", 40, 250) for c in conns] +
            [(tag, "fade", 250, 500)])


def s_workflow(prs):
    s = slide_of(prs)
    chrome = frame(s, "How It Works", "Only the consumer can close the case.")
    steps = [("1", "Share", "post + evidence", CORAL),
             ("2", "Discuss", "community · Me Too", PLUM),
             ("3", "Respond", "brand replies", TEAL),
             ("4", "Propose", "refund / repair / apology", GOLD),
             ("5", "Confirm", "only YOU close it", TEAL_DK)]
    y = 3.15
    circles, arrows = [], []
    for i, (n, t, d, c) in enumerate(steps):
        cx = 1.75 + i * 2.45
        cc = disc(s, cx, y, 0.62, fill=c)
        set_text(cc, n, 30, WHITE, bold=True, align=PP_ALIGN.CENTER, font=HEAD)
        textbox(s, cx - 1.05, y + 0.78, 2.1, 0.4, t, 15, INKTEXT, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
        textbox(s, cx - 1.1, y + 1.2, 2.2, 0.5, d, 10.5, MUTED, font=BODY, align=PP_ALIGN.CENTER)
        circles.append(cc)
        if i < 4:
            a = _shape(s, MSO_SHAPE.RIGHT_ARROW, cx + 0.72, y - 0.11, 0.95, 0.22, fill=LINE)
            arrows.append(a)
    band = rrect(s, 0.9, 5.35, 11.52, 1.35, fill=INK, radius=0.06)
    textbox(s, 1.3, 5.5, 10.9, 0.5, "Brands cannot unilaterally mark a complaint resolved.",
            17, GOLD, bold=True, font=HEAD)
    textbox(s, 1.3, 6.05, 10.9, 0.5,
            "The full, immutable resolution history stays public forever — that permanence is the trust moat.",
            12.5, TEAL_LT, font=BODY)
    pagechrome(s, 4)
    transition(s, "push")
    seq = []
    for i in range(5):
        seq.append((circles[i], "rise", 160 if i else 220, 340))
        if i < 4:
            seq.append((arrows[i], "wipe", 60, 220))
    animate(s, cascade(chrome, "fade", 100, 50, 300) + seq + [(band, "fade", 250, 600)])


def s_market(prs):
    s = slide_of(prs)
    chrome = frame(s, "Market Opportunity", "India-first — a large, underserved wedge.")
    bars = [("63M+", "Indian MSMEs with no formal helpdesk", 0.95, CORAL),
            ("₹ Bn", "Consumer-complaint & reputation search demand", 0.55, GOLD),
            ("$50B+", "Global customer-experience software TAM", 0.78, TEAL)]
    anim = []
    for i, (val, label, frac, c) in enumerate(bars):
        y = 2.35 + i * 1.15
        track = rrect(s, 0.95, y, 8.6, 0.72, fill=PAPER_2, radius=0.5)
        fill = rrect(s, 0.95, y, 8.6 * frac, 0.72, fill=c, radius=0.5)
        textbox(s, 1.2, y + 0.02, 6.5, 0.68, label, 12.5, WHITE, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, 9.7, y - 0.03, 2.6, 0.8, val, 26, INKTEXT, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
        anim += [track, fill]
    callout = rrect(s, 0.95, 6.0, 11.45, 0.72, fill=TEAL_LT, line=TEAL, line_w=1.5, radius=0.12)
    set_text(callout, "Wedge:  rank on “[Brand] complaints” search  +  SME SaaS priced far below Zendesk / Freshdesk",
             13, INKTEXT, bold=True, align=PP_ALIGN.CENTER, font=BODY)
    textbox(s, 9.7, 5.55, 2.7, 0.35, "illustrative", 9, MUTED, font=BODY, align=PP_ALIGN.RIGHT)
    pagechrome(s, 5)
    transition(s, "push")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + cascade(anim, "wipe", 200, 120, 500) +
            [(callout, "fade", 250, 500)])


def s_product(prs):
    s = slide_of(prs)
    chrome = frame(s, "Product — MVP is Live", "A full-stack platform, deployed and demo-ready.")
    feats = [("✍", "Complaint wizard", "AI draft + quality check", CORAL),
             ("🔎", "Smart search", "natural language + synonyms", TEAL),
             ("📊", "Brand dashboard", "KPIs + complaint inbox", GOLD),
             ("🔁", "Resolution loop", "propose → confirm", PLUM),
             ("📰", "Public feed", "trending + recently resolved", TEAL_DK),
             ("🎯", "Contextual ads", "AI targeting, never biased", CORAL),
             ("🏷", "B2B pages", "pricing + API docs", GOLD),
             ("🤖", "AI layer", "Azure GPT + rule fallbacks", PLUM)]
    tiles = []
    for i, (g, t, d, c) in enumerate(feats):
        col, row = i % 4, i // 4
        x = 0.9 + col * 2.98
        y = 2.15 + row * 1.55
        card = rrect(s, x, y, 2.72, 1.4, fill=WHITE, line=LINE, line_w=1.2, radius=0.1)
        icon_disc(s, x + 0.5, y + 0.42, 0.29, g, c)
        textbox(s, x + 0.9, y + 0.16, 1.75, 0.55, t, 12.5, INKTEXT, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE, leading=0.95)
        textbox(s, x + 0.24, y + 0.82, 2.3, 0.5, d, 10, MUTED, font=BODY)
        tiles.append(card)
    demo = rrect(s, 0.9, 5.42, 11.52, 1.15, fill=INK, radius=0.06)
    textbox(s, 1.25, 5.52, 3, 0.4, "LIVE DEMO", 12, GOLD, bold=True, font=BODY, tracking=2)
    textbox(s, 1.25, 5.9, 10.8, 0.55,
            "github.com/himsinghvi/panchayat    ·    login  admin / demo123    ·    8 personas · 6 brands · seeded cases",
            12.5, CREAM, font=BODY)
    pagechrome(s, 6)
    transition(s, "push")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + cascade(tiles, "rise", 180, 85, 360) +
            [(demo, "fade", 250, 600)])


def s_b2b(prs):
    s = slide_of(prs)
    chrome = frame(s, "For Business", "Turn a chaotic inbox into a case-management system.")
    # native "case ticket" mock (left)
    card = rrect(s, 0.95, 2.2, 5.1, 4.15, fill=WHITE, line=LINE, line_w=1.4, radius=0.05)
    rect(s, 0.95, 2.2, 5.1, 0.72, TEAL_DK)
    textbox(s, 1.25, 2.28, 3.5, 0.55, "Case #1042", 16, WHITE, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
    pill(s, 4.35, 2.36, 1.45, 0.4, "RESOLVED", GOLD, INK, size=10)
    rows = [("Customer", "Priya S. · Pune"), ("Issue", "AC install delayed 3 weeks"),
            ("Owner", "Vikram (support)"), ("SLA", "First reply 2h · closed 5d"),
            ("Channel", "Email → auto-synced")]
    for i, (k, v) in enumerate(rows):
        y = 3.15 + i * 0.56
        textbox(s, 1.25, y, 1.6, 0.4, k, 11, MUTED, bold=True, font=BODY)
        textbox(s, 2.75, y, 3.1, 0.4, v, 12, INKTEXT, font=BODY)
    prog = rrect(s, 1.25, 5.95, 4.5, 0.22, fill=PAPER_2, radius=0.5)
    rrect(s, 1.25, 5.95, 4.5, 0.22, fill=TEAL, radius=0.5)
    # feature list (right)
    feats = [("🌐", "Branded page  panchaayat.in/yourbrand", CORAL),
             ("✉", "Two-way email → case sync", TEAL),
             ("🆔", "Case ID, status & owner on every query", GOLD),
             ("🔒", "Private by default; public escalation if ignored", PLUM),
             ("📈", "CSAT / NPS + SLA dashboards", TEAL_DK),
             ("🎁", "14-day free trial · cancel anytime", CORAL)]
    fr = []
    for i, (g, t, c) in enumerate(feats):
        y = 2.28 + i * 0.7
        ic = icon_disc(s, 6.85, y + 0.24, 0.24, g, c)
        textbox(s, 7.25, y, 5.0, 0.5, t, 13, INKTEXT, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
        fr.append(ic)
    pagechrome(s, 7)
    transition(s, "push")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + [(card, "rise", 200, 550)] +
            cascade(fr, "fade", 200, 110, 350))


def s_pricing(prs):
    s = slide_of(prs)
    chrome = frame(s, "SaaS Pricing", "Simple, affordable plans for SMEs & MSMEs.")
    plans = [("Startup", "₹1,999", ["150 cases / month", "1 seat", "Branded page", "Email + SMS alerts"], False, TEAL),
             ("Scale-up", "₹3,999", ["Unlimited cases", "10 seats", "Email → case sync", "CSAT + social"], True, CORAL),
             ("Enterprise", "₹5,999", ["Full REST API", "White-label", "AI triage", "Unlimited seats"], False, PLUM)]
    cards = []
    for i, (name, price, feats, hot, c) in enumerate(plans):
        x = 1.55 + i * 3.5
        h = 4.25 if hot else 3.95
        y = 2.2 if hot else 2.35
        card = rrect(s, x, y, 3.0, h, fill=(WHITE if not hot else INK), line=(LINE if not hot else CORAL),
                     line_w=(1.25 if not hot else 2.5), radius=0.07)
        tcol = INKTEXT if not hot else CREAM
        mut = MUTED if not hot else TEAL_LT
        if hot:
            pill(s, x + 0.85, y - 0.2, 1.3, 0.4, "POPULAR", CORAL, WHITE, size=9.5)
        textbox(s, x + 0.3, y + 0.28, 2.4, 0.5, name, 17, c if not hot else CORAL, bold=True, font=HEAD)
        pt = textbox(s, x + 0.3, y + 0.8, 2.4, 0.7, price, 30, tcol, bold=True, font=HEAD)
        textbox(s, x + 0.32, y + 1.5, 2.4, 0.3, "per month", 10.5, mut, font=BODY)
        rect(s, x + 0.32, y + 1.92, 2.35, 0.02, (LINE if not hot else INK_SOFT))
        set_text(textbox(s, x + 0.32, y + 2.05, 2.45, 1.9, "", 11, tcol),
                 "\n".join("✓   " + f for f in feats), 11.5, tcol, font=BODY, leading=1.35, anchor=MSO_ANCHOR.TOP)
        cards.append(card)
    note = textbox(s, 0.9, 6.5, 11.5, 0.4, "No setup fee    ·    30-day money-back guarantee    ·    GST invoice",
                   11.5, MUTED, font=BODY, align=PP_ALIGN.CENTER, tracking=1)
    pagechrome(s, 8)
    transition(s, "push")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + cascade(cards, "rise", 200, 220, 500) +
            [(note, "fade", 200, 400)])


def s_revenue(prs):
    s = slide_of(prs)
    chrome = frame(s, "Business Model", "Trust-first monetization.")
    # native proportion bar
    segs = [("SaaS subscriptions", 0.55, TEAL), ("Enterprise API", 0.25, CORAL),
            ("Contextual ads", 0.12, GOLD), ("Analytics (future)", 0.08, PLUM)]
    x = 0.95; total_w = 11.45; y = 2.35
    anim = []
    for name, frac, c in segs:
        w = total_w * frac
        seg = rrect(s, x, y, w - 0.04, 0.9, fill=c, radius=0.06)
        set_text(seg, f"{int(frac*100)}%", 16, WHITE, bold=True, align=PP_ALIGN.CENTER, font=HEAD)
        x += w
        anim.append(seg)
    # legend rows
    lx = 0.95
    for i, (name, frac, c) in enumerate(segs):
        col = i % 2; rowi = i // 2
        gx = 0.95 + col * 6.0
        gy = 3.75 + rowi * 0.95
        disc(s, gx + 0.15, gy + 0.2, 0.15, fill=c)
        textbox(s, gx + 0.45, gy - 0.02, 5.3, 0.45, name, 14, INKTEXT, bold=True, font=BODY)
        sub = {"SaaS subscriptions": "primary revenue — Startup → Enterprise",
               "Enterprise API": "banks, insurers, telcos",
               "Contextual ads": "legal aid · warranty · refund tools",
               "Analytics (future)": "benchmarks · researcher API"}[name]
        textbox(s, gx + 0.45, gy + 0.35, 5.3, 0.4, sub, 10.5, MUTED, font=BODY)
    guard = rrect(s, 0.95, 5.85, 11.45, 0.75, fill=INK, radius=0.1)
    set_text(guard, "✗   We never let brands pay to remove or bury a legitimate complaint.",
             13.5, CORAL_LT, bold=True, align=PP_ALIGN.CENTER, font=BODY)
    pagechrome(s, 9)
    transition(s, "push")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + cascade(anim, "wipe", 200, 150, 450) +
            [(guard, "rise", 300, 500)])


def s_competition(prs):
    s = slide_of(prs)
    chrome = frame(s, "Competitive Landscape", "No incumbent owns the full resolution loop.")
    cols = ["Google\nReviews", "MouthShut", "NCH /\nE-Jagriti", "Zendesk", "Panchaayat"]
    feats = ["Public discussion", "Resolution loop", "Consumer confirms", "SME case mgmt", "India-first SEO"]
    grid = {(0, 0), (0, 1), (1, 1), (3, 3),
            (0, 4), (1, 4), (2, 4), (3, 4), (4, 4)}  # (row, col) that are checked
    x0, y0, cw, ch = 4.15, 2.55, 1.62, 0.72
    # column headers
    for j, cname in enumerate(cols):
        cx = x0 + j * cw
        hot = (j == 4)
        if hot:
            rrect(s, cx + 0.06, y0 - 0.62, cw - 0.12, 5 * ch + 0.75, fill=TEAL_LT).line.fill.background()
        set_text(textbox(s, cx, y0 - 0.62, cw, 0.55, "", 11, INKTEXT),
                 cname, 11, (TEAL_DK if hot else INKTEXT), bold=hot, align=PP_ALIGN.CENTER, font=BODY, leading=0.9)
    # rows
    for i, fname in enumerate(feats):
        ry = y0 + i * ch
        textbox(s, 0.95, ry, 3.1, ch, fname, 12.5, INKTEXT, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
        for j in range(5):
            cx = x0 + j * cw
            has = (i, j) in grid
            hot = (j == 4)
            mark = "✓" if has else "·"
            mc = (TEAL_DK if hot else MUTED) if has else LINE
            set_text(textbox(s, cx, ry, cw, ch, "", 16, mc),
                     mark, 17 if has else 15, mc, bold=has, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # hairline
        rect(s, 0.95, ry + ch - 0.01, 11.4, 0.012, LINE)
    pagechrome(s, 10)
    transition(s, "fade")
    animate(s, cascade(chrome, "fade", 100, 50, 300))


def s_gtm(prs):
    s = slide_of(prs)
    chrome = frame(s, "Go-to-Market", "Hyper-local density → national scale.")
    phases = [("Phase 1", "Seed Pune &\nHyderabad", CORAL),
              ("Phase 2", "SEO: “[brand]\ncomplaints”", GOLD),
              ("Phase 3", "SME outbound —\nD2C & appliances", TEAL),
              ("Phase 4", "Enterprise API —\nbanks & insurers", PLUM)]
    chevs = []
    for i, (p, d, c) in enumerate(phases):
        x = 0.95 + i * 2.95
        ch = _shape(s, MSO_SHAPE.CHEVRON, x, 2.25, 2.95, 1.75, fill=c)
        set_text(ch, f"{p}\n{d}", 12, WHITE, bold=True, align=PP_ALIGN.CENTER, font=BODY, leading=1.05)
        ch.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
        chevs.append(ch)
    rows = [("🔥", "Viral growth via Me Too, trending & recently-resolved feeds"),
            ("🎁", "Free trial converts SMEs; email-sync lock-in curbs churn"),
            ("🏛", "Roadmap: partner with NCH / E-Jagriti for escalations")]
    rr = []
    for i, (g, t) in enumerate(rows):
        y = 4.6 + i * 0.72
        ic = icon_disc(s, 1.3, y + 0.24, 0.26, g, TEAL)
        textbox(s, 1.75, y, 10.4, 0.5, t, 13.5, INKTEXT, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
        rr.append(ic)
    pagechrome(s, 11)
    transition(s, "push")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + cascade(chevs, "wipe", 200, 200, 420) +
            cascade(rr, "fade", 180, 120, 340))


def s_traction(prs):
    s = slide_of(prs)
    chrome = frame(s, "Traction & KPIs", "Where we are today — and what we measure.")
    tiles = [("MVP", "live & deployed", TEAL), ("8", "demo personas", CORAL),
             ("6", "seeded brands", GOLD), ("3", "SaaS tiers", PLUM)]
    ts = []
    for i, (big, lab, c) in enumerate(tiles):
        x = 0.95 + i * 2.9
        card = rrect(s, x, 2.2, 2.65, 1.55, fill=WHITE, line=c, line_w=2.0, radius=0.1)
        textbox(s, x, 2.32, 2.65, 0.85, big, 34, c, bold=True, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        textbox(s, x, 3.16, 2.65, 0.4, lab, 11, MUTED, font=BODY, align=PP_ALIGN.CENTER)
        ts.append(card)
    textbox(s, 0.95, 4.05, 11, 0.4, "North-star metrics we will report", 14, INKTEXT, bold=True, font=HEAD)
    kpis = [("⏱", "Brand response under 7 days", CORAL),
            ("✅", "Consumer-confirmed resolution rate", TEAL),
            ("📉", "Median time-to-resolution", GOLD),
            ("💳", "SME trial → paid conversion", PLUM)]
    ks = []
    for i, (g, t, c) in enumerate(kpis):
        y = 4.6 + i * 0.6
        ic = icon_disc(s, 1.28, y + 0.2, 0.22, g, c)
        textbox(s, 1.7, y - 0.03, 10.5, 0.45, t, 13, INKTEXT, bold=True, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
        ks.append(ic)
    pagechrome(s, 12)
    transition(s, "push")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + cascade(ts, "rise", 200, 150, 440) +
            cascade(ks, "fade", 150, 90, 300))


def s_roadmap(prs):
    s = slide_of(prs)
    chrome = frame(s, "Roadmap", "A 12–18 month product & GTM plan.")
    stops = [("Q2 ’26", "10 SME pilots\nPostgres + PWA", True),
             ("Q3 ’26", "Email→case\nHindi UI", False),
             ("Q4 ’26", "Social listening\n100 paying SMEs", False),
             ("Q1 ’27", "Enterprise API\n5 logos", False),
             ("Q2 ’27", "Mobile apps\nTier-2 cities", False)]
    y = 4.15
    line = rect(s, 1.3, y - 0.02, 10.7, 0.05, TEAL)
    nodes = []
    for i, (q, d, now) in enumerate(stops):
        cx = 1.7 + i * 2.5
        n = disc(s, cx, y, 0.16, fill=(CORAL if now else TEAL))
        above = (i % 2 == 0)
        cy = 2.45 if above else 4.75
        card = rrect(s, cx - 1.05, cy, 2.1, 1.1, fill=WHITE, line=LINE, line_w=1.25, radius=0.1)
        textbox(s, cx - 1.0, cy + 0.08, 2.0, 0.4, q, 13, (CORAL if now else TEAL_DK), bold=True, font=HEAD, align=PP_ALIGN.CENTER)
        textbox(s, cx - 1.0, cy + 0.48, 2.0, 0.6, d, 10.5, MUTED, font=BODY, align=PP_ALIGN.CENTER, leading=1.0)
        connector(s, cx, y, cx, cy + (1.1 if above else 0.0), color=LINE, w=1.25)
        nodes += [card, n]
    pagechrome(s, 13)
    transition(s, "fade")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + [(line, "wipe", 200, 500)] +
            cascade(nodes, "rise", 150, 110, 380))


def s_risks(prs):
    s = slide_of(prs)
    chrome = frame(s, "Risks & Mitigations", "An honest assessment for investors.")
    risks = [("Defamation", "Clear ToS · evidence requirement · takedown & appeal flow", RED),
             ("Fake reviews", "Verified-purchase signal + active moderation queue", CORAL),
             ("Cold start", "Hyper-local seeding + SEO complaint pages", GOLD),
             ("Brand apathy", "Public pages rank for “[brand] complaints” on Google", TEAL),
             ("SME churn", "Email-sync lock-in + 30-day money-back guarantee", PLUM)]
    rows = []
    for i, (r, m, c) in enumerate(risks):
        y = 2.2 + i * 0.88
        tag = rrect(s, 0.95, y, 2.75, 0.68, fill=c, radius=0.12)
        set_text(tag, r, 13.5, WHITE, bold=True, align=PP_ALIGN.CENTER, font=HEAD)
        arr = _shape(s, MSO_SHAPE.RIGHT_ARROW, 3.85, y + 0.22, 0.5, 0.24, fill=LINE)
        mc = rrect(s, 4.55, y, 7.85, 0.68, fill=WHITE, line=LINE, line_w=1.2, radius=0.1)
        set_text(mc, "   " + m, 12.5, INKTEXT, font=BODY, anchor=MSO_ANCHOR.MIDDLE)
        rows += [tag, arr, mc]
    pagechrome(s, 14)
    transition(s, "push")
    animate(s, cascade(chrome, "fade", 100, 50, 300) + cascade(rows, "fade", 200, 90, 340))


def s_ask(prs):
    s = slide_of(prs)
    bg(s, INK)
    disc(s, 12.0, 1.0, 2.3, fill=None, line=INK_SOFT, line_w=2.0)
    em = emblem(s, 1.05, 1.0, 0.3, ring=TEAL, center=CORAL, dots=CREAM)
    k = textbox(s, 1.5, 0.72, 8, 0.4, "THE ASK", 13, TEAL, bold=True, font=BODY, tracking=3)
    t = textbox(s, 0.85, 1.35, 11, 0.9, "Partner with us to build the trust layer.", 30, CREAM, bold=True, font=HEAD)
    blocks = [("₹ [X] Cr seed round", "Customize this figure before presenting.", CORAL),
              ("Use of funds", "40% GTM & sales   ·   30% engineering   ·   20% ops & legal   ·   10% buffer", TEAL),
              ("12-month milestones", "100 paying SMEs      ·      10K MAU      ·      1 enterprise pilot", GOLD)]
    cards = []
    for i, (h, d, c) in enumerate(blocks):
        y = 2.7 + i * 1.35
        card = rrect(s, 0.85, y, 11.6, 1.15, fill=INK_SOFT, line=c, line_w=2.0, radius=0.08)
        rect(s, 0.85, y, 0.12, 1.15, c)
        textbox(s, 1.25, y + 0.15, 11, 0.5, h, 19, c, bold=True, font=HEAD)
        textbox(s, 1.25, y + 0.66, 11, 0.4, d, 13, CREAM, font=BODY)
        cards.append(card)
    pagechrome(s, 15, dark=True)
    transition(s, "cover")
    animate(s, cascade(em, "fade", 120, 50, 250) + [(k, "fade", 150, 400), (t, "rise", 120, 550)] +
            cascade(cards, "rise", 200, 260, 520))


def s_close(prs):
    s = slide_of(prs)
    bg(s, INK)
    disc(s, 1.9, 6.4, 2.4, fill=None, line=INK_SOFT, line_w=2.0)
    disc(s, 1.9, 6.4, 1.7, fill=None, line=INK_SOFT, line_w=2.0)
    em = emblem(s, 6.66, 2.1, 0.5, ring=TEAL, center=CORAL, dots=CREAM)
    ty = textbox(s, 0.7, 2.95, 12, 1.2, "Thank you.", 60, CREAM, bold=True, font=HEAD, align=PP_ALIGN.CENTER)
    line = rect(s, 6.16, 4.25, 1.0, 0.06, CORAL)
    tag = textbox(s, 0.7, 4.45, 12, 0.5, "Make it visible.   Get it resolved.", 20, TEAL_LT, bold=True,
                  font=HEAD, align=PP_ALIGN.CENTER)
    contact = textbox(s, 0.7, 5.35, 12, 1.4,
                      "Demo    [your-vercel-url].vercel.app\n"
                      "GitHub    github.com/himsinghvi/panchayat\n"
                      "Email    [founder@email.com]",
                      14, CREAM, font=BODY, align=PP_ALIGN.CENTER, leading=1.5)
    transition(s, "fade")
    animate(s, cascade(em, "fade", 150, 60, 300) +
            [(ty, "rise", 150, 650), (line, "wipe", 150, 400), (tag, "fade", 200, 500),
             (contact, "fade", 250, 600)])


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for fn in (s_title, s_problem, s_solution, s_workflow, s_market, s_product, s_b2b,
               s_pricing, s_revenue, s_competition, s_gtm, s_traction, s_roadmap,
               s_risks, s_ask, s_close):
        fn(prs)
    prs.save(OUTPUT)
    print(f"Created: {OUTPUT} ({len(prs.slides)} slides, 16:9, native vector, animated)")


if __name__ == "__main__":
    build()
